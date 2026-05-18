"""Next Steps slide — W11/D2 (mode-agnostic).

Numbered list. Each step has:
  - action verb + body
  - owner role (or fallback)
  - timing (or fallback)
  - optional dependencies

Pulls from ``payload.next_steps`` which can be either a list[str]
(common case) or a list[dict] with ``{action, owner_role, timing,
dependencies}`` keys. Falls back to a 3-bullet "Recommended next
actions" derived from the recommendation prose if next_steps is
empty.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list, get_recommendation_text
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide


def _extract_steps(payload: Any) -> list[dict[str, Any]]:
    raw = _coerce_to_list(payload_get(payload, "next_steps", default=[]))
    out: list[dict[str, Any]] = []
    for item in raw:
        if isinstance(item, dict):
            action = str(item.get("action") or item.get("text") or "").strip()
            if not action:
                continue
            out.append({
                "action": action,
                "owner_role": str(item.get("owner_role") or item.get("owner") or "").strip(),
                "timing": str(item.get("timing") or item.get("timeline") or "").strip(),
                "dependencies": _coerce_to_list(item.get("dependencies") or []),
            })
        elif isinstance(item, str) and item.strip():
            out.append({"action": item.strip(), "owner_role": "", "timing": "", "dependencies": []})
    return out


def _fallback_steps(payload: Any) -> list[dict[str, Any]]:
    """Synthesize 3 bullets from the recommendation prose when
    next_steps is missing entirely — better than an empty slide."""
    rec = get_recommendation_text(payload)
    if not rec:
        return []
    # Split on sentence-end punctuation, take up to 3 non-trivial pieces.
    parts: list[str] = []
    for chunk in rec.replace("\n", " ").split("."):
        s = chunk.strip()
        if len(s) > 6:
            parts.append(s + ".")
        if len(parts) >= 3:
            break
    return [{"action": p, "owner_role": "", "timing": "", "dependencies": []} for p in parts]


@register_slide("next_steps")
class NextStepsSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        add_horizontal_band(
            slide, left=0.0, top=0.0, width=SLIDE_WIDTH_IN, height=0.4,
            color_hex=str(primary),
        )
        add_textbox(
            slide, left=0.5, top=0.5, width=SLIDE_WIDTH_IN - 1.0, height=0.5,
            text="Next Steps",
            font_size=24, bold=True,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        steps = _extract_steps(payload) or _fallback_steps(payload)
        steps = steps[:7]  # cap so the slide stays readable

        body = add_textbox(
            slide,
            left=0.5, top=1.3,
            width=SLIDE_WIDTH_IN - 1.0, height=5.5,
            text="",
            font_size=12,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        if not steps:
            add_paragraph(
                body.text_frame,
                "(no next steps recorded by the writer)",
                font_size=12,
                color=parse_hex(DEFAULT_MUTED),
            )
        else:
            for i, step in enumerate(steps, start=1):
                # Main line: "1. Action"
                add_paragraph(
                    body.text_frame,
                    f"{i}. {step['action'][:220]}",
                    font_size=13, bold=True,
                    color=parse_hex(secondary_hex),
                )
                # Owner / timing / deps as a smaller subline if any are set.
                meta_bits: list[str] = []
                if step["owner_role"]:
                    meta_bits.append(f"Owner: {step['owner_role']}")
                if step["timing"]:
                    meta_bits.append(f"Timing: {step['timing']}")
                if step["dependencies"]:
                    meta_bits.append(
                        "Deps: " + ", ".join(str(d) for d in step["dependencies"][:3])
                    )
                if meta_bits:
                    add_paragraph(
                        body.text_frame,
                        "    " + " · ".join(meta_bits),
                        font_size=10,
                        color=parse_hex(DEFAULT_MUTED),
                    )

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])
