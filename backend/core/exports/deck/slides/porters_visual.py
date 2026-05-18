"""Porter's Five Forces visual slide — W11/D3.

Canonical Porter's layout: a central Rivalry box flanked by four
surrounding force boxes (New Entrants on top, Substitutes on bottom,
Suppliers on left, Buyers on right). Each box carries:

  - Force name (bold header)
  - Intensity badge (red/amber/green pill)
  - Rationale (2 lines)
  - Key drivers (2 bullets max)
  - Citation chip (top-right corner)

Below the diagram: "Overall attractiveness: {rating} — {rationale}".

Reads ``payload.frameworks.porters_five_forces``. If absent (W8 Run B
writer-truncation case), renders the documented fallback.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .._layout import (
    DEFAULT_FONT,
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_textbox,
    parse_hex,
)
from ..shape_helpers import (
    add_citation_chip,
    add_intensity_badge,
    add_text_in_box,
)
from ..._base import payload_get
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide

# Box geometry. Slide is 13.333 × 7.5 inches; central rivalry sits
# in the middle and the four surrounding forces orbit it.
_BOX_W = 3.6
_BOX_H = 2.0
_CENTER_X = SLIDE_WIDTH_IN / 2
_CENTER_Y = 4.2  # slightly below true vertical centre so title + footer fit


def _force_position(force: str) -> tuple[float, float]:
    """Return (left, top) for the given force, all in inches."""
    half_w = _BOX_W / 2
    half_h = _BOX_H / 2
    if force == "rivalry":
        return (_CENTER_X - half_w, _CENTER_Y - half_h)
    if force == "new_entrant_threat":
        return (_CENTER_X - half_w, _CENTER_Y - half_h - _BOX_H - 0.25)
    if force == "substitute_threat":
        return (_CENTER_X - half_w, _CENTER_Y - half_h + _BOX_H + 0.25)
    if force == "supplier_power":
        return (_CENTER_X - half_w - _BOX_W - 0.25, _CENTER_Y - half_h)
    if force == "buyer_power":
        return (_CENTER_X - half_w + _BOX_W + 0.25, _CENTER_Y - half_h)
    return (0.5, 0.5)


_FORCE_LABELS: dict[str, str] = {
    "rivalry": "Competitive Rivalry",
    "new_entrant_threat": "Threat of New Entrants",
    "substitute_threat": "Threat of Substitutes",
    "supplier_power": "Supplier Power",
    "buyer_power": "Buyer Power",
}


def _draw_force_box(
    slide: Any,
    *,
    force: str, force_data: dict[str, Any],
    primary_hex: str, secondary_hex: str,
    cited: list[str],
) -> None:
    left, top = _force_position(force)
    label = _FORCE_LABELS.get(force, force.replace("_", " ").title())

    # Outer rectangle — coloured border, transparent fill so the
    # force-name header line reads cleanly.
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(_BOX_W), Inches(_BOX_H),
    )
    rect.fill.background()
    rect.line.color.rgb = parse_hex(primary_hex)
    rect.line.width = Pt(1.0)

    # Force-name header strip (filled with primary, white text).
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(_BOX_W), Inches(0.32),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = parse_hex(primary_hex)
    header.line.fill.background()
    tf = header.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = parse_hex("#FFFFFF")

    # Intensity badge — top-right of header strip.
    intensity = str(force_data.get("intensity") or "").lower().strip()
    add_intensity_badge(
        slide,
        left=left + _BOX_W - 0.75,
        top=top + 0.04,
        width=0.65, height=0.24,
        intensity=intensity,
    )

    # Rationale + key drivers below the header.
    body_top = top + 0.38
    body_height = _BOX_H - 0.42
    body_left = left + 0.1
    body_width = _BOX_W - 0.2

    rationale = str(force_data.get("rationale") or "").strip()
    key_drivers = [str(d).strip() for d in (force_data.get("key_drivers") or []) if str(d).strip()]
    citations = [str(c).strip() for c in (force_data.get("evidence_citations") or []) if str(c).strip()]

    # Rationale takes ~55% of body height.
    rationale_h = body_height * 0.55
    add_text_in_box(
        slide,
        left=body_left, top=body_top,
        width=body_width, height=rationale_h,
        text=rationale or "(no rationale captured)",
        font_size=9,
        color=parse_hex(DEFAULT_MUTED if not rationale else secondary_hex),
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )
    # Key drivers — bottom band.
    drivers_top = body_top + rationale_h + 0.02
    drivers_h = body_height - rationale_h - 0.02
    drivers_text = "  ·  ".join(key_drivers[:2]) if key_drivers else ""
    if drivers_text:
        add_text_in_box(
            slide,
            left=body_left, top=drivers_top,
            width=body_width, height=drivers_h,
            text=drivers_text,
            font_size=9, bold=True,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

    # Citation chip — top-right corner, on top of the header strip.
    if citations:
        cid = citations[0]
        if cid not in cited:
            cited.append(cid)
        add_citation_chip(
            slide,
            left=left + _BOX_W - 0.34,
            top=top + 0.36,
            number=len(cited),
            claim_id=cid,
            primary_hex=primary_hex,
        )


@register_slide("porters_five_forces_visual")
class PortersVisualSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        add_horizontal_band(
            slide, left=0.0, top=0.0, width=SLIDE_WIDTH_IN, height=0.4,
            color_hex=str(primary),
        )

        frameworks = payload_get(payload, "frameworks", default={}) or {}
        p5 = frameworks.get("porters_five_forces") if isinstance(frameworks, dict) else None
        if not isinstance(p5, dict):
            add_textbox(
                slide, left=0.5, top=0.5, width=SLIDE_WIDTH_IN - 1.0, height=0.5,
                text="Industry Forces — Porter's Five Forces",
                font_size=24, bold=True,
                color=parse_hex(secondary_hex),
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            )
            add_textbox(
                slide, left=0.5, top=1.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
                text="Porter's Five Forces — not produced for this engagement.",
                font_size=12,
                color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )
            return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])

        # Title
        market = str(p5.get("market_definition") or "").strip()
        title_text = (
            f"Industry Forces — {market}" if market else "Industry Forces — Porter's Five Forces"
        )
        # Truncate aggressively so the long-form market definitions
        # don't blow the title row.
        if len(title_text) > 110:
            title_text = title_text[:107].rstrip() + "…"
        add_textbox(
            slide, left=0.5, top=0.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
            text=title_text,
            font_size=22, bold=True,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        cited: list[str] = []
        for force in (
            "new_entrant_threat", "substitute_threat",
            "supplier_power", "buyer_power", "rivalry",
        ):
            fd = p5.get(force) or {}
            if isinstance(fd, dict):
                _draw_force_box(
                    slide,
                    force=force, force_data=fd,
                    primary_hex=str(primary), secondary_hex=secondary_hex,
                    cited=cited,
                )

        # Bottom strip: overall attractiveness + rationale.
        overall = str(p5.get("overall_attractiveness") or "").lower().strip()
        overall_rationale = str(p5.get("overall_rationale") or "").strip()
        bottom_top = 7.0
        if overall:
            add_intensity_badge(
                slide,
                left=0.5, top=bottom_top + 0.03,
                width=1.2, height=0.28,
                intensity=overall,
            )
            add_text_in_box(
                slide,
                left=1.85, top=bottom_top,
                width=SLIDE_WIDTH_IN - 2.5, height=0.4,
                text=f"Overall attractiveness — {overall_rationale[:240]}"
                     if overall_rationale
                     else f"Overall attractiveness rated {overall.upper()}.",
                font_size=10,
                color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
            )

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=cited)
