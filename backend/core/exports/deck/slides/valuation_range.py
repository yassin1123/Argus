"""Valuation Range slide — W11/D2 (M&A-specific).

Three coloured boxes side-by-side: Low | Base | High.
Each box shows:
  - £m value (large)
  - methodology (small text below)
  - up to 3 key_assumptions bullets

Bottom strip: comparable_transactions_cited summarized as a
1-line "X comparable transactions cited (range A.B–C.Dx multiple)"
hint, since the slide can't fit 10 deal rows.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide


_SCENARIO_COLOURS: dict[str, str] = {
    "low":  "#5b6470",
    "base": "#0F6E56",
    "high": "#B8860B",
}


def _scenario_block(
    vr: dict[str, Any], key: str
) -> tuple[float | None, str, list[str]]:
    """Pull ``(gbp_m, methodology, key_assumptions)`` from one of the
    valuation scenarios, tolerant of both nested ({"low":{"gbp_m":...}})
    and flat ({"low_gbp_m":...}) shapes."""
    node = vr.get(key) or {}
    if isinstance(node, dict) and "gbp_m" in node:
        try:
            gbp = float(node["gbp_m"])
        except (TypeError, ValueError):
            gbp = None
        method = str(node.get("methodology") or "").strip()
        assumps = _coerce_to_list(node.get("key_assumptions") or [])
        return gbp, method, [str(x).strip() for x in assumps if str(x).strip()][:3]
    flat = vr.get(f"{key}_gbp_m")
    if isinstance(flat, (int, float)):
        return float(flat), "", []
    return None, "", []


@register_slide("valuation_range")
class ValuationRangeSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        add_horizontal_band(
            slide, left=0.0, top=0.0, width=SLIDE_WIDTH_IN, height=0.4,
            color_hex=str(primary),
        )
        add_textbox(
            slide, left=0.5, top=0.5, width=SLIDE_WIDTH_IN - 1.0, height=0.5,
            text="Valuation Range",
            font_size=24, bold=True,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        vr = payload_get(payload, "valuation_range", default=None)
        if not isinstance(vr, dict):
            add_textbox(
                slide, left=0.5, top=1.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
                text="Valuation range not produced for this engagement.",
                font_size=12, color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )
            return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])

        # Three boxes — equal width, distinct colours.
        box_top = 1.4
        box_height = 4.6
        margin_x = 0.5
        gap = 0.3
        box_w = (SLIDE_WIDTH_IN - 2 * margin_x - 2 * gap) / 3

        for i, key in enumerate(("low", "base", "high")):
            left = margin_x + i * (box_w + gap)
            colour_hex = _SCENARIO_COLOURS[key]

            # Coloured header strip on top of the box.
            add_horizontal_band(
                slide, left=left, top=box_top,
                width=box_w, height=0.45,
                color_hex=colour_hex,
            )
            add_textbox(
                slide, left=left, top=box_top + 0.04,
                width=box_w, height=0.42,
                text=key.upper(),
                font_size=14, bold=True,
                color=parse_hex("#FFFFFF"),
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

            gbp, method, assumps = _scenario_block(vr, key)
            value_text = f"£{gbp:.1f}m" if gbp is not None else "—"
            add_textbox(
                slide, left=left, top=box_top + 0.6,
                width=box_w, height=1.0,
                text=value_text,
                font_size=32, bold=True,
                color=parse_hex(colour_hex),
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            add_textbox(
                slide, left=left + 0.1, top=box_top + 1.7,
                width=box_w - 0.2, height=0.7,
                text=method[:80] if method else "(methodology not specified)",
                font_size=10,
                color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
            )
            # Assumptions bullets.
            assump_box = add_textbox(
                slide, left=left + 0.1, top=box_top + 2.5,
                width=box_w - 0.2, height=box_height - 2.6,
                text="", font_size=10,
                color=parse_hex(secondary_hex),
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            )
            for a in assumps:
                add_paragraph(
                    assump_box.text_frame, a[:140],
                    font_size=10, bullet=True,
                    color=parse_hex(secondary_hex),
                )
            if not assumps:
                add_paragraph(
                    assump_box.text_frame, "(no anchoring assumptions provided)",
                    font_size=9, color=parse_hex(DEFAULT_MUTED),
                )

        # Bottom strip: comparable transactions summary.
        comps = _coerce_to_list(vr.get("comparable_transactions_cited") or [])
        n_comps = sum(1 for c in comps if isinstance(c, dict))
        if n_comps:
            multiples = [
                str(c.get("multiple") or "").strip()
                for c in comps if isinstance(c, dict)
            ]
            multiples = [m for m in multiples if m]
            mult_blurb = ", ".join(multiples[:4])
            comp_text = (
                f"{n_comps} comparable transaction(s) cited"
                + (f" — {mult_blurb}" if mult_blurb else "")
            )
        else:
            comp_text = "(no comparable transactions cited)"
        add_textbox(
            slide, left=0.5, top=box_top + box_height + 0.15,
            width=SLIDE_WIDTH_IN - 1.0, height=0.4,
            text=comp_text,
            font_size=11,
            color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])
