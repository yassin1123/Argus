"""Integration Plan slide — W11/D2 (M&A-specific).

Three-band horizontal layout:
  Band 1 — Day 1 priorities       (list[str])
  Band 2 — First 100 days          (list[InitiativeBlock])
  Band 3 — First year              (list[InitiativeBlock])

Each InitiativeBlock surfaces as a bulleted line:
    workstream — owner_role — milestone

Bottom-right badge: integration_complexity_rating (low/medium/high)
keyed to the same severity palette as the risks-matrix slide.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide

_COMPLEXITY_COLOURS: dict[str, str] = {
    "low": "#0F6E56",
    "medium": "#B8860B",
    "high": "#B91C1C",
}


def _initiative_line(block: Any) -> tuple[str, list[str]]:
    """Return ``(line_text, citation_ids)`` for one InitiativeBlock-ish
    dict (or string)."""
    if isinstance(block, dict):
        ws = str(block.get("workstream") or "").strip()
        owner = str(block.get("owner_role") or "").strip()
        ms = str(block.get("milestone") or "").strip()
        line_parts = [p for p in (ws, owner) if p]
        prefix = " · ".join(line_parts)
        line = f"{prefix} — {ms}" if prefix and ms else (prefix or ms or "—")
        cids = [str(c) for c in (block.get("evidence_citations") or []) if c]
        return line, cids
    return str(block).strip(), []


def _render_band(
    slide: Any,
    *, title: str, items: list[Any],
    left: float, top: float, width: float, height: float,
    primary_hex: str, secondary_hex: str,
    cited: list[str],
) -> None:
    add_horizontal_band(
        slide, left=left, top=top, width=width, height=0.4,
        color_hex=primary_hex,
    )
    add_textbox(
        slide, left=left + 0.15, top=top + 0.04,
        width=width - 0.3, height=0.35,
        text=title,
        font_size=12, bold=True,
        color=parse_hex("#FFFFFF"),
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    body = add_textbox(
        slide, left=left + 0.15, top=top + 0.5,
        width=width - 0.3, height=height - 0.55,
        text="", font_size=10,
        color=parse_hex(secondary_hex),
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )
    if not items:
        add_paragraph(
            body.text_frame, "(no items)",
            font_size=10, color=parse_hex(DEFAULT_MUTED),
        )
        return
    for item in items[:5]:
        line, cids = _initiative_line(item)
        if not line:
            continue
        add_paragraph(
            body.text_frame, line[:220],
            font_size=10, bullet=True,
            color=parse_hex(secondary_hex),
        )
        for c in cids:
            if c not in cited:
                cited.append(c)


@register_slide("integration_plan")
class IntegrationPlanSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
        deck_context: Any = None,
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        ip = payload_get(payload, "integration_plan", default=None)
        if not isinstance(ip, dict):
            add_textbox(
                slide, left=0.5, top=1.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
                text="Integration plan not produced for this engagement.",
                font_size=12, color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )
            return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])

        cited: list[str] = []
        band_top = 1.3
        band_height = 1.75
        band_left = 0.5
        band_width = SLIDE_WIDTH_IN - 1.0

        _render_band(
            slide, title="Day 1 priorities",
            items=_coerce_to_list(ip.get("day_one_priorities") or []),
            left=band_left, top=band_top,
            width=band_width, height=band_height,
            primary_hex=str(primary), secondary_hex=secondary_hex,
            cited=cited,
        )
        _render_band(
            slide, title="First 100 days",
            items=_coerce_to_list(ip.get("first_100_days") or []),
            left=band_left, top=band_top + band_height + 0.2,
            width=band_width, height=band_height,
            primary_hex=str(primary), secondary_hex=secondary_hex,
            cited=cited,
        )
        _render_band(
            slide, title="First year",
            items=_coerce_to_list(ip.get("first_year") or []),
            left=band_left, top=band_top + 2 * (band_height + 0.2),
            width=band_width, height=band_height,
            primary_hex=str(primary), secondary_hex=secondary_hex,
            cited=cited,
        )

        # Complexity badge bottom-right.
        rating = str(ip.get("integration_complexity_rating") or "").lower().strip()
        if rating in _COMPLEXITY_COLOURS:
            badge_w, badge_h = 1.4, 0.45
            badge_left = SLIDE_WIDTH_IN - badge_w - 0.5
            badge_top = 7.5 - badge_h - 0.25
            add_horizontal_band(
                slide, left=badge_left, top=badge_top,
                width=badge_w, height=badge_h,
                color_hex=_COMPLEXITY_COLOURS[rating],
            )
            add_textbox(
                slide, left=badge_left, top=badge_top + 0.04,
                width=badge_w, height=badge_h - 0.08,
                text=f"COMPLEXITY: {rating.upper()}",
                font_size=10, bold=True,
                color=parse_hex("#FFFFFF"),
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

        return SlideResult(
            slide_index=len(presentation.slides) - 1,
            citation_ids=cited,
        )
