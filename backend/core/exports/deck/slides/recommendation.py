"""Recommendation slide — W11/D1.

Layout:
  - Header band + "Recommendation" title.
  - Upper half: full recommendation prose (color-coded by verdict).
  - Lower half left: "Conditions" / "Decision criteria" (mode-aware:
    M&A pulls deal_structure_implications.walk_away_triggers; everyone
    else pulls kill_criteria + decision_criteria).
  - Lower half right: source-panel summary ("3 SEC filings · 1
    transcript · 12 firm-library chunks") aggregated from payload.sources.
"""

from __future__ import annotations

from collections import Counter
from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
    verdict_colour,
)
from ..._base import payload_get
from ...one_pager_renderer import (
    _coerce_to_list,
    _label_source_type,
    classify_recommendation,
    get_recommendation_text,
)
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide


def _walk_away_triggers(payload: Any) -> list[str]:
    ds = payload_get(payload, "deal_structure_implications", default={}) or {}
    if isinstance(ds, dict):
        out = _coerce_to_list(ds.get("walk_away_triggers") or [])
        return [str(x).strip() for x in out if str(x).strip()]
    return []


def _decision_criteria(payload: Any) -> list[str]:
    """Generic decision-criteria fallback for non-M&A modes."""
    kc = _coerce_to_list(payload_get(payload, "kill_criteria", default=[]))
    dc = _coerce_to_list(payload_get(payload, "decision_criteria", default=[]))
    items: list[str] = []
    for x in list(kc) + list(dc):
        if isinstance(x, dict):
            s = str(x.get("text") or x.get("description") or "").strip()
        else:
            s = str(x).strip()
        if s:
            items.append(s)
    return items


def _source_panel_summary(payload: Any) -> list[str]:
    """Aggregate sources[].type into 'N label' strings, sorted by count."""
    srcs = _coerce_to_list(payload_get(payload, "sources", default=[]))
    types: list[str] = []
    for s in srcs:
        if isinstance(s, dict):
            t = s.get("type") or s.get("source_type") or ""
        else:
            t = str(s)
        types.append(t)
    if not types:
        return []
    counter: Counter[str] = Counter()
    for t in types:
        counter[_label_source_type(t)] += 1
    return [
        f"{n} {label}"
        for label, n in sorted(counter.items(), key=lambda kv: (-kv[1], kv[0]))
    ]


@register_slide("recommendation")
class RecommendationSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
        deck_context: Any = None,
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)

        # ---- Upper half: full recommendation prose. ----
        rec_text = get_recommendation_text(payload) or "(no recommendation produced)"
        rec_color = classify_recommendation(rec_text)
        add_textbox(
            slide,
            left=0.5, top=1.3,
            width=SLIDE_WIDTH_IN - 1.0, height=2.3,
            text=rec_text[:1200],  # safety bound — single slide cannot scroll
            font_size=16, bold=True,
            color=verdict_colour(rec_color),
            align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
        )

        # ---- Lower half left: conditions / decision criteria ----
        triggers = _walk_away_triggers(payload)
        if triggers:
            block_title = "Walk-away triggers"
            items = triggers[:5]
        else:
            block_title = "Decision criteria"
            items = _decision_criteria(payload)[:5]

        add_textbox(
            slide,
            left=0.5, top=3.9,
            width=8.0, height=0.45,
            text=block_title,
            font_size=14, bold=True,
            color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )
        criteria_shape = add_textbox(
            slide,
            left=0.5, top=4.4,
            width=8.0, height=2.7,
            text="",
            font_size=12,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT,
        )
        for item in items:
            add_paragraph(
                criteria_shape.text_frame,
                str(item)[:200],
                font_size=12, bullet=True,
                color=parse_hex(secondary_hex),
            )
        if not items:
            add_paragraph(
                criteria_shape.text_frame,
                "(none provided by the writer)",
                font_size=11, color=parse_hex(DEFAULT_MUTED),
            )

        # ---- Lower half right: source panel ----
        add_textbox(
            slide,
            left=9.0, top=3.9,
            width=SLIDE_WIDTH_IN - 9.5, height=0.45,
            text="Sources",
            font_size=14, bold=True,
            color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )
        sources_shape = add_textbox(
            slide,
            left=9.0, top=4.4,
            width=SLIDE_WIDTH_IN - 9.5, height=2.7,
            text="",
            font_size=11,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT,
        )
        rows = _source_panel_summary(payload)
        for row in rows[:8]:
            add_paragraph(
                sources_shape.text_frame, row,
                font_size=11, bullet=True,
                color=parse_hex(secondary_hex),
            )
        if not rows:
            add_paragraph(
                sources_shape.text_frame, "(no sources recorded)",
                font_size=10, color=parse_hex(DEFAULT_MUTED),
            )

        # Recommendation claim_ids cited from this slide.
        rec_claim_ids = _coerce_to_list(
            payload_get(payload, "recommendation_claim_ids", default=[])
        )
        return SlideResult(
            slide_index=len(presentation.slides) - 1,
            citation_ids=[str(c) for c in rec_claim_ids if c],
        )
