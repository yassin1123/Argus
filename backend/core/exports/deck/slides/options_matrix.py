"""Strategic Options slide — W11/D2 (growth-specific, text-stub).

Per spec hard rule: Day 2 is a text-only options listing. Day 3
replaces this builder body with a real 2x2 visual via the W8
TwoByTwoMatrix framework.

Data sources, in order:
  1. ``payload.options_matrix`` (list of dicts with name/quadrant/rationale).
  2. ``payload.frameworks.two_by_two.items`` (W8 schema).
  3. ``payload.key_reasons[:4]`` synthesized into option-like lines.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list, _stringify_item
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide


def _extract_options(payload: Any) -> list[dict[str, Any]]:
    """Pull a list of ``{name, rationale, quadrant?}`` dicts."""
    direct = _coerce_to_list(payload_get(payload, "options_matrix", default=[]))
    out: list[dict[str, Any]] = []
    for item in direct:
        if not isinstance(item, dict):
            continue
        name = str(item.get("name") or item.get("option") or "").strip()
        if not name:
            continue
        out.append({
            "name": name,
            "rationale": str(item.get("rationale") or item.get("description") or "").strip(),
            "quadrant": str(item.get("quadrant") or "").strip(),
            "citations": list(item.get("evidence_citations") or []),
        })
    if out:
        return out

    fw = payload_get(payload, "frameworks", default={}) or {}
    if isinstance(fw, dict):
        tb = fw.get("two_by_two") or {}
        if isinstance(tb, dict):
            for it in _coerce_to_list(tb.get("items") or []):
                if not isinstance(it, dict):
                    continue
                name = str(it.get("name") or "").strip()
                if not name:
                    continue
                out.append({
                    "name": name,
                    "rationale": str(it.get("rationale") or "").strip(),
                    "quadrant": str(it.get("quadrant") or "").strip(),
                    "citations": list(it.get("evidence_citations") or []),
                })
    if out:
        return out

    # Last-resort fallback: synthesize lines from key_reasons.
    for r in _coerce_to_list(payload_get(payload, "key_reasons", default=[]))[:4]:
        s = _stringify_item(r)
        if s:
            out.append({"name": s[:80], "rationale": s, "quadrant": "", "citations": []})
    return out


@register_slide("options_matrix")
class OptionsMatrixSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        add_horizontal_band(
            slide, left=0.0, top=0.0, width=SLIDE_WIDTH_IN, height=0.4,
            color_hex=str(primary),
        )
        add_textbox(
            slide, left=0.5, top=0.5, width=SLIDE_WIDTH_IN - 1.0, height=0.5,
            text="Strategic Options",
            font_size=24, bold=True,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        add_textbox(
            slide, left=0.5, top=1.0, width=SLIDE_WIDTH_IN - 1.0, height=0.4,
            text="(Day 3 turns this into a 2×2 visual; today it lists the options as text.)",
            font_size=10, color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )

        options = _extract_options(payload)
        body = add_textbox(
            slide, left=0.5, top=1.5,
            width=SLIDE_WIDTH_IN - 1.0, height=5.5,
            text="", font_size=12,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        cited: list[str] = []
        if not options:
            add_paragraph(
                body.text_frame,
                "(no strategic options surfaced by the writer)",
                font_size=11, color=parse_hex(DEFAULT_MUTED),
            )
        else:
            for opt in options[:6]:
                add_paragraph(
                    body.text_frame,
                    opt["name"][:160],
                    font_size=14, bold=True,
                    color=parse_hex(secondary_hex),
                )
                if opt["rationale"]:
                    add_paragraph(
                        body.text_frame,
                        f"    {opt['rationale'][:240]}",
                        font_size=11,
                        color=parse_hex(secondary_hex),
                    )
                if opt["quadrant"]:
                    add_paragraph(
                        body.text_frame,
                        f"    quadrant: {opt['quadrant']}",
                        font_size=10,
                        color=parse_hex(DEFAULT_MUTED),
                    )
                for c in opt["citations"]:
                    cs = str(c)
                    if cs and cs not in cited:
                        cited.append(cs)

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=cited)
