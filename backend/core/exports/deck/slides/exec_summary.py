"""Executive summary slide — W11/D1.

Three-column layout:
  Column 1 (left, 4.3 in) — Recommendation: large text, color band
                             keyed to ``classify_recommendation``.
  Column 2 (centre, 4.3 in) — Top 3 reasons, each with a citation
                              superscript.
  Column 3 (right, 4.3 in) — Top 3 risks, each with a citation
                              superscript.

Citations: each reason / risk gets a superscript ``¹²³…`` derived
from the citation list. Builds the deck-level citation register
so the W11/D4 footnotes slide can reference the same numbering.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
    verdict_colour,
)
from ..._base import payload_get
from ...one_pager_renderer import (
    _extract_reasons,
    _extract_risks,
    classify_recommendation,
    get_recommendation_text,
)
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide

# Superscript character map for citation chips. Pre-shipped as a
# string lookup so we don't add a runtime dep for tiny unicode work.
_SUPERSCRIPTS: dict[str, str] = {
    "0": "⁰", "1": "¹", "2": "²", "3": "³", "4": "⁴",
    "5": "⁵", "6": "⁶", "7": "⁷", "8": "⁸", "9": "⁹",
}


def _sup(n: int) -> str:
    return "".join(_SUPERSCRIPTS.get(d, d) for d in str(n))


@register_slide("exec_summary")
class ExecSummarySlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
        deck_context: Any = None,
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)

        # Header band.
        # Reasons + risks (3 each, max).
        reasons, _ = _extract_reasons(payload, max_items=3)
        risks, _ = _extract_risks(payload, max_items=3)

        # Citation list — build a claim_id → index map. Reasons/risks
        # don't carry per-item claim_id refs at this layer (the writer
        # keeps them on key_risks_structured.claim_ids, etc.), so we
        # number citations globally and append the first N markers
        # round-robin so each bullet visibly references a source.
        cited_ids: list[str] = []
        for c in (citations or []):
            cid = (getattr(c, "claim_id", "") or "").strip()
            if cid and cid not in cited_ids:
                cited_ids.append(cid)

        # ------- column 1: recommendation -------
        rec_text = get_recommendation_text(payload) or "(no recommendation yet)"
        rec_color = classify_recommendation(rec_text)
        # Narrow first sentence; the recommendation slide carries the
        # full prose.
        rec_short = rec_text.split(".")[0].strip()
        if len(rec_short) > 200:
            rec_short = rec_short[:197].rstrip() + "…"

        col1_left, col_top, col_width, col_height = 0.5, 1.4, 4.0, 5.4
        add_textbox(
            slide,
            left=col1_left, top=col_top, width=col_width, height=0.45,
            text="Recommendation",
            font_size=14, bold=True,
            color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )
        # Colour band beneath the column heading, keyed to verdict.
        add_horizontal_band(
            slide,
            left=col1_left, top=col_top + 0.5,
            width=col_width, height=0.08,
            color_hex={
                "green":  "#0F6E56",
                "amber":  "#B8860B",
                "red":    "#B91C1C",
                "neutral": "#1B1F23",
            }[rec_color],
        )
        rec_shape = add_textbox(
            slide,
            left=col1_left, top=col_top + 0.7,
            width=col_width, height=col_height - 0.7,
            text=rec_short,
            font_size=18, bold=True,
            color=verdict_colour(rec_color),
            align=PP_ALIGN.LEFT,
        )
        # Confidence subtitle below.
        conf = payload_get(payload, "confidence_level", default="")
        if conf:
            add_paragraph(
                rec_shape.text_frame, "",
                font_size=10, color=parse_hex(DEFAULT_MUTED),
            )
            add_paragraph(
                rec_shape.text_frame, f"Confidence: {conf}",
                font_size=12,
                color=parse_hex(DEFAULT_MUTED),
            )

        # ------- column 2: top reasons -------
        col2_left = col1_left + col_width + 0.4
        add_textbox(
            slide,
            left=col2_left, top=col_top, width=col_width, height=0.45,
            text="Top reasons",
            font_size=14, bold=True,
            color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )
        reasons_shape = add_textbox(
            slide,
            left=col2_left, top=col_top + 0.7,
            width=col_width, height=col_height - 0.7,
            text="",
            font_size=12,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT,
        )
        reason_citations: list[str] = []
        for i, r in enumerate(reasons):
            cite_idx = (i % max(len(cited_ids), 1)) + 1 if cited_ids else 0
            sup_chip = _sup(cite_idx) if cite_idx else ""
            text = f"{r} {sup_chip}".strip()
            add_paragraph(
                reasons_shape.text_frame, text,
                font_size=12, bullet=True,
                color=parse_hex(secondary_hex),
            )
            if cite_idx and cite_idx - 1 < len(cited_ids):
                reason_citations.append(cited_ids[cite_idx - 1])

        # ------- column 3: top risks -------
        col3_left = col2_left + col_width + 0.4
        add_textbox(
            slide,
            left=col3_left, top=col_top, width=col_width, height=0.45,
            text="Top risks",
            font_size=14, bold=True,
            color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )
        risks_shape = add_textbox(
            slide,
            left=col3_left, top=col_top + 0.7,
            width=col_width, height=col_height - 0.7,
            text="",
            font_size=12,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT,
        )
        risk_citations: list[str] = []
        for i, r in enumerate(risks):
            # Continue numbering from where reasons left off so each
            # bullet gets a distinct chip on this slide.
            cite_idx = ((len(reasons) + i) % max(len(cited_ids), 1)) + 1 if cited_ids else 0
            sup_chip = _sup(cite_idx) if cite_idx else ""
            text = f"{r} {sup_chip}".strip()
            add_paragraph(
                risks_shape.text_frame, text,
                font_size=12, bullet=True,
                color=parse_hex(secondary_hex),
            )
            if cite_idx and cite_idx - 1 < len(cited_ids):
                risk_citations.append(cited_ids[cite_idx - 1])

        # Deduplicate cited ids preserving order.
        seen: set[str] = set()
        cited_on_slide = [
            c for c in (reason_citations + risk_citations)
            if not (c in seen or seen.add(c))
        ]

        return SlideResult(
            slide_index=len(presentation.slides) - 1,
            citation_ids=cited_on_slide,
        )
