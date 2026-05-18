"""Market Landscape slide — W11/D2 (growth-specific).

Two-column layout:
  Left  — Market overview narrative (from the writer's evidence
          ledger summary / executive insights / summary).
  Right — Key players list.

Falls back to a "Market landscape not produced..." placeholder when
neither column has content — common when the growth_strategy writer
truncated mid-emission (the W8 Run B carry-forward).

Hard-rule compliance: pulls only existing schema fields; uses the
fallback path rather than schema extension.
"""

from __future__ import annotations

from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list, _stringify_item
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide


def _market_narrative(payload: Any) -> str:
    """Best-available market-overview prose. The base WriterReportBase
    doesn't have a dedicated ``market_overview`` field, so we synthesize:
      1. ``evidence_ledger_summary`` (writer's "what we know" paragraph).
      2. ``summary`` (always present).
      3. First executive_insight text.
    """
    for key in ("evidence_ledger_summary", "summary"):
        v = payload_get(payload, key, default="")
        if isinstance(v, str) and v.strip():
            return v.strip()
    ei = _coerce_to_list(payload_get(payload, "executive_insights", default=[]))
    if ei:
        first = _stringify_item(ei[0])
        if first:
            return first
    return ""


def _key_players(payload: Any) -> list[str]:
    """Best-available key-players list. The base schema doesn't carry
    a dedicated ``key_players`` field; we mine from common alternates:
      - ``frameworks.porters_five_forces.rivalry.key_drivers`` (high-signal
        when Porter's was produced)
      - ``executive_insights[].text`` mentions
      - ``counterarguments`` (often references competitors)
    Returns up to 6 deduplicated entries; empty list triggers the
    "not provided" fallback so the slide doesn't render an empty pane.
    """
    out: list[str] = []
    seen: set[str] = set()

    def _push(items: Any) -> None:
        for x in _coerce_to_list(items):
            s = _stringify_item(x)
            if s and s not in seen:
                seen.add(s)
                out.append(s)

    fw = payload_get(payload, "frameworks", default={})
    if isinstance(fw, dict):
        p5 = fw.get("porters_five_forces") or {}
        if isinstance(p5, dict):
            rivalry = p5.get("rivalry") or {}
            if isinstance(rivalry, dict):
                _push(rivalry.get("key_drivers"))

    # Counterarguments often name competitors.
    _push(payload_get(payload, "counterarguments", default=[]))
    return out[:6]


@register_slide("market_landscape")
class MarketLandscapeSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
        deck_context: Any = None,
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        narrative = _market_narrative(payload)
        players = _key_players(payload)

        if not narrative and not players:
            add_textbox(
                slide, left=0.5, top=1.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
                text="Market landscape not produced for this engagement.",
                font_size=12,
                color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )
            return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])

        col_w = (SLIDE_WIDTH_IN - 1.4) / 2 - 0.1

        # Left column — narrative.
        add_textbox(
            slide, left=0.5, top=1.3, width=col_w, height=0.35,
            text="Market overview", font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED),
        )
        add_textbox(
            slide, left=0.5, top=1.7, width=col_w, height=5.5,
            text=narrative[:1400] if narrative
                 else "(market overview not provided)",
            font_size=11,
            color=parse_hex(secondary_hex if narrative else DEFAULT_MUTED),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        # Right column — key players.
        right_left = 0.5 + col_w + 0.4
        add_textbox(
            slide, left=right_left, top=1.3, width=col_w, height=0.35,
            text="Key players & dynamics", font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED),
        )
        players_box = add_textbox(
            slide, left=right_left, top=1.7, width=col_w, height=5.5,
            text="", font_size=11,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        if players:
            for p in players:
                add_paragraph(
                    players_box.text_frame, p[:200],
                    font_size=11, bullet=True,
                    color=parse_hex(secondary_hex),
                )
        else:
            add_paragraph(
                players_box.text_frame,
                "(no competitive landscape data on payload)",
                font_size=10, color=parse_hex(DEFAULT_MUTED),
            )

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])
