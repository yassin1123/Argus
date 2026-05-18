"""Financial Profile slide — W11/D2 (M&A-specific).

Three-row layout:
  Top: Revenue trajectory chart (column chart, 4-5 periods)
  Middle: EBITDA trajectory + margin profile (text)
  Bottom: working_capital + debt + capex + cash_flow_quality (3 bullets)

The chart uses python-pptx's native ``add_chart`` with
``XL_CHART_TYPE.COLUMN_CLUSTERED``. If the trajectory has fewer than
2 data points (writer cut content short) we fall back to a textual
listing — no broken chart frame.

Per spec hard rule: don't pull payload fields that don't exist on
the base schema. We read ``financial_profile`` defensively and
short-circuit to a "not produced" placeholder when absent rather
than crashing the deck.
"""

from __future__ import annotations

from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .._layout import (
    DEFAULT_FONT,
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide


def _trajectory_points(fp: dict[str, Any], key: str) -> list[tuple[str, float]]:
    """Pull (period, value_gbp_m) pairs from
    ``revenue_trajectory.points`` (or ``ebitda_trajectory.points``).
    Coerces tolerantly: skip points missing either field, cast value
    to float, keep order.
    """
    traj = fp.get(key) or {}
    if not isinstance(traj, dict):
        return []
    pts = _coerce_to_list(traj.get("points") or [])
    out: list[tuple[str, float]] = []
    for p in pts:
        if not isinstance(p, dict):
            continue
        period = str(p.get("period") or "").strip()
        v = p.get("value_gbp_m")
        try:
            v_f = float(v) if v is not None else None
        except (TypeError, ValueError):
            v_f = None
        if period and v_f is not None:
            out.append((period, v_f))
    return out


def _add_revenue_chart(
    slide: Any, points: list[tuple[str, float]],
    *, left: float, top: float, width: float, height: float,
    primary_hex: str,
) -> None:
    cats = [p[0] for p in points]
    values = [p[1] for p in points]
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Revenue (£m)", values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        data,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    # Colour the data series to match firm primary.
    try:
        plot = chart.plots[0]
        ser = plot.series[0]
        fill = ser.format.fill
        fill.solid()
        fill.fore_color.rgb = parse_hex(primary_hex)
    except Exception:
        # Some pptx versions don't expose .format on a column series;
        # the default chart styling is acceptable as fallback.
        pass


@register_slide("financial_profile")
class FinancialProfileSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
        deck_context: Any = None,
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        fp = payload_get(payload, "financial_profile", default=None)
        if not isinstance(fp, dict):
            add_textbox(
                slide, left=0.5, top=1.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
                text="Financial profile not produced for this engagement.",
                font_size=12,
                color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )
            return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])

        # Top row: revenue trajectory.
        rev_points = _trajectory_points(fp, "revenue_trajectory")
        chart_left, chart_top, chart_w, chart_h = 0.5, 1.3, 6.5, 3.0
        add_textbox(
            slide, left=chart_left, top=chart_top, width=chart_w, height=0.35,
            text="Revenue trajectory (£m)",
            font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED), align=PP_ALIGN.LEFT,
        )
        if len(rev_points) >= 2:
            _add_revenue_chart(
                slide, rev_points,
                left=chart_left, top=chart_top + 0.45,
                width=chart_w, height=chart_h - 0.45,
                primary_hex=str(primary),
            )
        else:
            add_textbox(
                slide, left=chart_left, top=chart_top + 0.5,
                width=chart_w, height=chart_h - 0.5,
                text="(revenue trajectory has fewer than 2 data points)",
                font_size=11, color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )

        # Top-right: margin profile + EBITDA snapshot.
        right_left = chart_left + chart_w + 0.4
        right_w = SLIDE_WIDTH_IN - right_left - 0.5
        add_textbox(
            slide, left=right_left, top=chart_top, width=right_w, height=0.35,
            text="Margin profile",
            font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED), align=PP_ALIGN.LEFT,
        )
        mp = fp.get("margin_profile") or {}
        margin_box = add_textbox(
            slide, left=right_left, top=chart_top + 0.45,
            width=right_w, height=chart_h - 0.45,
            text="",
            font_size=11, color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        if isinstance(mp, dict) and any(mp.get(k) for k in ("gross_margin", "ebitda_margin", "fcf_margin")):
            for label, key in (("Gross margin", "gross_margin"),
                                ("EBITDA margin", "ebitda_margin"),
                                ("FCF margin", "fcf_margin")):
                v = str(mp.get(key) or "").strip()
                if v:
                    add_paragraph(
                        margin_box.text_frame,
                        f"{label}: {v}",
                        font_size=11, bullet=True,
                        color=parse_hex(secondary_hex),
                    )
            trend = str(mp.get("trend_commentary") or "").strip()
            if trend:
                add_paragraph(margin_box.text_frame, "", font_size=6)
                add_paragraph(
                    margin_box.text_frame, trend[:240],
                    font_size=10, color=parse_hex(DEFAULT_MUTED),
                )
        else:
            add_paragraph(
                margin_box.text_frame, "(margin profile not provided)",
                font_size=10, color=parse_hex(DEFAULT_MUTED),
            )

        # Bottom row: WC / debt / capex / cash-flow quality bullets.
        add_textbox(
            slide, left=0.5, top=4.6, width=SLIDE_WIDTH_IN - 1.0, height=0.35,
            text="Capital + cash flow",
            font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED), align=PP_ALIGN.LEFT,
        )
        bottom_box = add_textbox(
            slide, left=0.5, top=5.0, width=SLIDE_WIDTH_IN - 1.0, height=2.2,
            text="", font_size=11, color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        for label, key in (
            ("Working capital", "working_capital_dynamics"),
            ("Debt", "debt_structure"),
            ("Capex", "capex_intensity"),
            ("Cash flow quality", "cash_flow_quality"),
        ):
            v = str(fp.get(key) or "").strip()
            if v:
                add_paragraph(
                    bottom_box.text_frame,
                    f"{label}: {v[:280]}",
                    font_size=11, bullet=True,
                    color=parse_hex(secondary_hex),
                )

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])
