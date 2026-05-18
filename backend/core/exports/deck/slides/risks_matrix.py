"""Risks & Mitigations slide — W11/D2 (mode-agnostic).

Table layout: Risk | Severity | Mitigation columns.

  - M&A mode: rows come from payload.risks_and_mitigations
    (list[RiskAssessment]) with real severity / mitigation / residual_risk.
  - Other modes: rows synthesized from payload.risks (list[str]) +
    key_risks_structured[].text if present; severity inferred as
    "medium" by default; mitigation populated from a residual blurb.

Severity colour-coded per cell (red / amber / green).
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .._layout import (
    DEFAULT_FONT,
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list, _stringify_item
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide

_SEVERITY_COLOURS: dict[str, str] = {
    "high": "#B91C1C",
    "medium": "#B8860B",
    "low": "#0F6E56",
}

_DEFAULT_SEVERITY = "medium"
_MAX_ROWS = 6  # 1 header + 5 risk rows; keeps the table on one A4 slide.


def _extract_rows(payload: Any) -> list[dict[str, Any]]:
    """Normalize to a list of ``{description, severity, mitigation,
    claim_ids}`` rows from whatever payload shape we have.

    M&A path:  ``risks_and_mitigations`` list of RiskAssessment dicts.
    Other:     ``risks`` (list[str]) merged with
               ``key_risks_structured[].text/claim_ids``.
    """
    rms = _coerce_to_list(payload_get(payload, "risks_and_mitigations", default=[]))
    out: list[dict[str, Any]] = []
    for r in rms:
        if not isinstance(r, dict):
            continue
        desc = str(r.get("description") or r.get("text") or "").strip()
        if not desc:
            continue
        sev = str(r.get("severity") or _DEFAULT_SEVERITY).lower().strip()
        if sev not in _SEVERITY_COLOURS:
            sev = _DEFAULT_SEVERITY
        out.append({
            "description": desc,
            "severity": sev,
            "mitigation": str(r.get("mitigation") or "").strip()
                          or "—",
            "claim_ids": list(r.get("evidence_citations") or r.get("claim_ids") or []),
        })
        if len(out) >= (_MAX_ROWS - 1):
            return out
    if out:
        return out

    # Non-M&A fallback: flat risks[] + key_risks_structured[]
    flat = [_stringify_item(x) for x in
            _coerce_to_list(payload_get(payload, "risks", default=[]))]
    krs = _coerce_to_list(payload_get(payload, "key_risks_structured", default=[]))
    for i, desc in enumerate([x for x in flat if x][: _MAX_ROWS - 1]):
        # Pair with key_risks_structured by index when possible.
        claim_ids: list[str] = []
        mitigation = "—"
        if i < len(krs) and isinstance(krs[i], dict):
            claim_ids = list(krs[i].get("claim_ids") or [])
            mitigation = str(
                krs[i].get("mitigation") or krs[i].get("response") or ""
            ).strip() or "—"
        out.append({
            "description": desc,
            "severity": _DEFAULT_SEVERITY,
            "mitigation": mitigation,
            "claim_ids": [str(c) for c in claim_ids if c],
        })
    return out


def _set_cell_text(cell: Any, text: str, *, font_size: int = 10, bold: bool = False,
                   color: RGBColor | None = None, align: int = PP_ALIGN.LEFT,
                   fill: RGBColor | None = None) -> None:
    """python-pptx cell-styling helper. Cells default to empty text, so we
    fully replace via the cell's text_frame paragraph runs."""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    # Clear any default run; pptx tables start with an empty run.
    for run in list(p.runs):
        run.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


@register_slide("risks_matrix")
class RisksMatrixSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        add_horizontal_band(
            slide, left=0.0, top=0.0, width=SLIDE_WIDTH_IN, height=0.4,
            color_hex=str(primary),
        )
        add_textbox(
            slide, left=0.5, top=0.5, width=SLIDE_WIDTH_IN - 1.0, height=0.5,
            text="Risks & Mitigations",
            font_size=24, bold=True,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        rows = _extract_rows(payload)
        if not rows:
            add_textbox(
                slide, left=0.5, top=1.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
                text="(no risks recorded for this engagement)",
                font_size=12,
                color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )
            return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])

        # Table dimensions.
        n_rows = 1 + len(rows)  # header + body
        n_cols = 3
        table_left = 0.5
        table_top = 1.3
        table_width = SLIDE_WIDTH_IN - 1.0
        # Header 0.45in, body rows 0.85in each.
        row_heights = [Inches(0.45)] + [Inches(0.85)] * len(rows)
        col_widths = [Inches(table_width * 0.45), Inches(table_width * 0.12), Inches(table_width * 0.43)]

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(table_left), Inches(table_top),
            Inches(table_width), sum(row_heights, Inches(0)),
        )
        tbl = table_shape.table
        for i, h in enumerate(row_heights):
            tbl.rows[i].height = h
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = w

        # Header row.
        header_bg = parse_hex(primary)
        for j, label in enumerate(["Risk", "Severity", "Mitigation"]):
            _set_cell_text(
                tbl.cell(0, j), label,
                font_size=11, bold=True,
                color=parse_hex("#FFFFFF"),
                align=PP_ALIGN.LEFT,
                fill=header_bg,
            )

        # Body rows.
        cited: list[str] = []
        for i, row in enumerate(rows, start=1):
            _set_cell_text(
                tbl.cell(i, 0),
                row["description"][:240],
                font_size=9,
                color=parse_hex(secondary_hex),
            )
            sev = row["severity"]
            _set_cell_text(
                tbl.cell(i, 1),
                sev.upper(),
                font_size=10, bold=True,
                color=parse_hex("#FFFFFF"),
                align=PP_ALIGN.CENTER,
                fill=parse_hex(_SEVERITY_COLOURS[sev]),
            )
            _set_cell_text(
                tbl.cell(i, 2),
                row["mitigation"][:300],
                font_size=9,
                color=parse_hex(secondary_hex),
            )
            for cid in row["claim_ids"]:
                if cid and cid not in cited:
                    cited.append(cid)

        return SlideResult(
            slide_index=len(presentation.slides) - 1,
            citation_ids=cited,
        )
