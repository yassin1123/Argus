"""Title slide — W11/D1 + D4.

Layout:
  - Top primary-coloured band (full width × 0.45 in).
  - Top-right corner: firm logo image (from the asset cache) OR
    firm-name text styled with the primary colour when the logo
    URL isn't reachable.
  - Centred recommendation prose as the title text in firm primary
    colour (large) — falls back to engagement title if recommendation
    is blank.
  - Subtitle: target name + month/year, in firm secondary colour.
  - Bottom-right: "Confidential. Prepared by {firm.name}." in
    secondary colour (the spec's branded subtitle line).

Per spec hard rule:
  - Logo fetch uses the asset cache (24h TTL on disk). The slide
    builder is sync so we drive the async fetch via ``asyncio.run`` —
    safe here because the DeckPptxExporter's ``render`` is called
    from the export service's awaited generate_artifact, not from a
    running event loop in the slide builder's stack.
    Actually safer: we never call .run inside an existing loop; the
    cache lookup is sync (read file from disk), and we resolve the
    URL → cached-PNG-bytes via a sync wrapper that uses ``asyncio.run``
    only if no event loop is currently running.

The title slide opts out of chrome (title bar + footer + footnotes
applied by DeckBuilder.finalize_chrome) via SlideResult.skip_chrome
because the title slide owns its own visual identity.
"""

from __future__ import annotations

import asyncio
import io
import logging
from datetime import datetime, timezone
from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...asset_cache import fetch_and_cache_logo
from ...one_pager_renderer import get_recommendation_text
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide

logger = logging.getLogger(__name__)


def _resolve_logo_sync(firm_id: Any, logo_url: str) -> bytes | None:
    """Sync wrapper around the async asset cache. Tolerant of being
    called inside or outside a running event loop.

    - Outside a loop: ``asyncio.run`` drives the coroutine.
    - Inside a loop (rare for the deck renderer but possible in tests):
      fall back to None — the title slide renders the firm-name
      placeholder rather than blocking. Tests that need to exercise
      the cached-bytes path can pre-populate the cache directory.
    """
    if not logo_url:
        return None
    try:
        loop = asyncio.get_running_loop()
    except RuntimeError:
        loop = None
    if loop is not None and loop.is_running():
        # We can't block on the async call from inside a running loop;
        # the cache is the only fast path. Read the cache file directly.
        from ...asset_cache import _cache_fresh, _cache_path
        p = _cache_path(firm_id)
        if _cache_fresh(p):
            try:
                return p.read_bytes()
            except OSError:
                return None
        return None
    try:
        return asyncio.run(fetch_and_cache_logo(firm_id, logo_url))
    except Exception as e:  # noqa: BLE001
        logger.info("logo resolve failed (%s) — using firm-name fallback", e)
        return None


@register_slide("title")
class TitleSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
        deck_context: Any = None,
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY
        firm_name = (
            (firm_branding or {}).get("_firm_name")
            or payload_get(payload, "_firm_name", default="Argus")
            or "Argus"
        )
        logo_url = str((firm_branding or {}).get("logo_url") or "").strip()
        firm_id = (firm_branding or {}).get("_firm_id") or firm_name

        slide = add_blank_slide(presentation)

        # Top brand band.
        add_horizontal_band(
            slide,
            left=0.0, top=0.0,
            width=SLIDE_WIDTH_IN, height=0.45,
            color_hex=str(primary),
        )

        # Top-right: firm logo (image) or firm-name text fallback.
        logo_bytes = _resolve_logo_sync(firm_id, logo_url)
        logo_h_in = 0.9
        logo_max_w_in = 2.5
        if logo_bytes:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(logo_bytes),
                    Inches(SLIDE_WIDTH_IN - logo_max_w_in - 0.5),
                    Inches(0.7),
                    height=Inches(logo_h_in),
                )
            except Exception as e:  # noqa: BLE001
                logger.info("logo embed failed (%s) — using firm-name fallback", e)
                logo_bytes = None
        if not logo_bytes:
            add_textbox(
                slide,
                left=SLIDE_WIDTH_IN - logo_max_w_in - 0.5, top=0.7,
                width=logo_max_w_in, height=0.7,
                text=firm_name,
                font_size=20, bold=True,
                color=parse_hex(primary),
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
            )

        # Recommendation as the title.
        rec = get_recommendation_text(payload) or str(
            payload_get(payload, "_engagement_title", default="Argus engagement")
        )
        rec_short = rec.split(".")[0].strip()
        if len(rec_short) > 110:
            rec_short = rec_short[:107].rstrip() + "…"

        add_textbox(
            slide,
            left=0.7, top=2.4,
            width=SLIDE_WIDTH_IN - 1.4, height=2.0,
            text=rec_short,
            font_size=36, bold=True,
            color=parse_hex(primary),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Subtitle row.
        target = (
            payload_get(payload, "_target_name", default="")
            or payload_get(payload, "_engagement_title", default="")
        )
        date_label = datetime.now(tz=timezone.utc).strftime("%B %Y")
        subtitle = " · ".join(p for p in (str(target), date_label) if p)
        add_textbox(
            slide,
            left=0.7, top=4.4,
            width=SLIDE_WIDTH_IN - 1.4, height=0.6,
            text=subtitle,
            font_size=16,
            color=parse_hex(secondary),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        # Bottom-right: "Confidential. Prepared by {firm name}."
        add_textbox(
            slide,
            left=0.7,
            top=SLIDE_HEIGHT_IN - 0.7,
            width=SLIDE_WIDTH_IN - 1.4, height=0.4,
            text=f"Confidential. Prepared by {firm_name}.",
            font_size=11,
            color=parse_hex(secondary),
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        return SlideResult(
            slide_index=len(presentation.slides) - 1,
            citation_ids=[],
            skip_chrome=True,
        )
