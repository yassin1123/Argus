"""Target Overview slide — W11/D2 (M&A-specific).

4-quadrant layout:
  TL: Business model paragraph
  TR: Segments table (name + revenue% + growth_rate)
  BL: Geographic exposure list
  BR: Ownership history + customer concentration

Skips gracefully when ``target_overview`` is absent (general-mode
fallback prints a one-line placeholder instead of failing the
whole deck render).
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .._layout import (
    DEFAULT_FONT,
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide


def _set_cell_text(cell: Any, text: str, *, font_size: int = 10, bold: bool = False,
                   color: RGBColor | None = None, align: int = PP_ALIGN.LEFT,
                   fill: RGBColor | None = None) -> None:
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = align
    for run in list(p.runs):
        run.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def _render_segments_table(
    slide: Any, segments: list[dict[str, Any]],
    left: float, top: float, width: float,
    *, primary_hex: str, secondary_hex: str,
) -> None:
    rows = segments[:6]
    if not rows:
        add_textbox(
            slide, left=left, top=top + 0.4, width=width, height=0.5,
            text="(no segment breakdown)",
            font_size=10, color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )
        return
    n_rows = 1 + len(rows)
    row_h = Inches(0.32)
    table_h = Inches(0.4) + row_h * n_rows
    tbl_shape = slide.shapes.add_table(
        n_rows, 3, Inches(left), Inches(top + 0.4),
        Inches(width), table_h,
    )
    tbl = tbl_shape.table
    for i in range(n_rows):
        tbl.rows[i].height = row_h
    tbl.columns[0].width = Inches(width * 0.55)
    tbl.columns[1].width = Inches(width * 0.20)
    tbl.columns[2].width = Inches(width * 0.25)
    header_bg = parse_hex(primary_hex)
    for j, label in enumerate(["Segment", "Revenue %", "Growth"]):
        _set_cell_text(
            tbl.cell(0, j), label,
            font_size=9, bold=True,
            color=parse_hex("#FFFFFF"),
            align=PP_ALIGN.LEFT,
            fill=header_bg,
        )
    for i, r in enumerate(rows, start=1):
        _set_cell_text(tbl.cell(i, 0), str(r.get("name") or "")[:80],
                       font_size=9, color=parse_hex(secondary_hex))
        pct = r.get("revenue_pct")
        pct_text = f"{pct:.1f}%" if isinstance(pct, (int, float)) else "—"
        _set_cell_text(tbl.cell(i, 1), pct_text, font_size=9,
                       color=parse_hex(secondary_hex), align=PP_ALIGN.RIGHT)
        _set_cell_text(tbl.cell(i, 2), str(r.get("growth_rate") or "—")[:30],
                       font_size=9, color=parse_hex(secondary_hex))


@register_slide("target_overview")
class TargetOverviewSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        add_horizontal_band(
            slide, left=0.0, top=0.0, width=SLIDE_WIDTH_IN, height=0.4,
            color_hex=str(primary),
        )
        add_textbox(
            slide, left=0.5, top=0.5, width=SLIDE_WIDTH_IN - 1.0, height=0.5,
            text="Target Overview",
            font_size=24, bold=True,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        to = payload_get(payload, "target_overview", default=None)
        if not isinstance(to, dict):
            add_textbox(
                slide, left=0.5, top=1.5, width=SLIDE_WIDTH_IN - 1.0, height=0.6,
                text="Target overview not produced for this engagement.",
                font_size=12,
                color=parse_hex(DEFAULT_MUTED),
                align=PP_ALIGN.LEFT,
            )
            return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])

        # Layout coordinates: a 2x2 grid inside the body region
        # (top 1.3 in, height 5.5 in).
        col_w = (SLIDE_WIDTH_IN - 1.4) / 2 - 0.1
        row_top_top = 1.3
        row_bot_top = 4.2
        col_left_l = 0.5
        col_left_r = 0.5 + col_w + 0.2

        # TL — Business model
        add_textbox(
            slide, left=col_left_l, top=row_top_top, width=col_w, height=0.35,
            text="Business model", font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED), align=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, left=col_left_l, top=row_top_top + 0.4,
            width=col_w, height=2.4,
            text=str(to.get("business_model") or "—")[:600],
            font_size=10,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        # TR — Segments
        add_textbox(
            slide, left=col_left_r, top=row_top_top, width=col_w, height=0.35,
            text="Segments", font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED), align=PP_ALIGN.LEFT,
        )
        segments = _coerce_to_list(to.get("segments") or [])
        _render_segments_table(
            slide, [s for s in segments if isinstance(s, dict)],
            left=col_left_r, top=row_top_top, width=col_w,
            primary_hex=str(primary), secondary_hex=secondary_hex,
        )

        # BL — Geographic exposure
        add_textbox(
            slide, left=col_left_l, top=row_bot_top, width=col_w, height=0.35,
            text="Geographic exposure", font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED), align=PP_ALIGN.LEFT,
        )
        geo_box = add_textbox(
            slide, left=col_left_l, top=row_bot_top + 0.4,
            width=col_w, height=2.4,
            text="",
            font_size=10,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        geos = _coerce_to_list(to.get("geographies") or [])
        if geos:
            for g in geos[:8]:
                if isinstance(g, dict):
                    geo = str(g.get("geography") or "")
                    pct = g.get("revenue_pct")
                    pct_text = f"{pct:.1f}%" if isinstance(pct, (int, float)) else ""
                    line = f"{geo} — {pct_text}" if pct_text else geo
                    if line.strip():
                        add_paragraph(
                            geo_box.text_frame, line,
                            font_size=10, bullet=True,
                            color=parse_hex(secondary_hex),
                        )
        else:
            add_paragraph(
                geo_box.text_frame, "(not provided)",
                font_size=10, color=parse_hex(DEFAULT_MUTED),
            )

        # BR — Ownership + customer concentration
        add_textbox(
            slide, left=col_left_r, top=row_bot_top, width=col_w, height=0.35,
            text="Ownership & customers", font_size=12, bold=True,
            color=parse_hex(DEFAULT_MUTED), align=PP_ALIGN.LEFT,
        )
        own_box = add_textbox(
            slide, left=col_left_r, top=row_bot_top + 0.4,
            width=col_w, height=2.4,
            text="",
            font_size=10,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )
        own = str(to.get("ownership_history") or "").strip()
        conc = str(to.get("key_customers_concentration") or "").strip()
        if own:
            add_paragraph(own_box.text_frame, "Ownership", font_size=10, bold=True,
                          color=parse_hex(DEFAULT_MUTED))
            add_paragraph(own_box.text_frame, own[:300], font_size=10,
                          color=parse_hex(secondary_hex))
        if conc:
            add_paragraph(own_box.text_frame, "", font_size=8)
            add_paragraph(own_box.text_frame, "Customer concentration", font_size=10, bold=True,
                          color=parse_hex(DEFAULT_MUTED))
            add_paragraph(own_box.text_frame, conc[:300], font_size=10,
                          color=parse_hex(secondary_hex))
        if not own and not conc:
            add_paragraph(own_box.text_frame, "(not provided)",
                          font_size=10, color=parse_hex(DEFAULT_MUTED))

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])
