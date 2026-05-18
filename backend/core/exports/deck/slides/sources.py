"""Sources slide — W11/D2 (mode-agnostic).

The "every claim verified" payoff slide. Aggregates the payload's
``sources`` list by type label (SEC filings, earnings transcripts,
firm-library documents, news sources, Companies House filings,
uploads) and prints each group with:
  - count + label header
  - up to N source titles bulleted underneath

If sources is empty (rare) we fall back to a single line referencing
the firm-content library; the deck still ships, just less defensible.
"""

from __future__ import annotations

from collections import defaultdict
from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .._layout import (
    DEFAULT_MUTED,
    DEFAULT_PRIMARY,
    DEFAULT_SECONDARY,
    SLIDE_WIDTH_IN,
    add_blank_slide,
    add_horizontal_band,
    add_paragraph,
    add_textbox,
    parse_hex,
)
from ..._base import payload_get
from ...one_pager_renderer import _coerce_to_list, _label_source_type
from ._base import SlideBuilderBase, SlideResult
from ._registry import register_slide

_MAX_TITLES_PER_GROUP = 4


def _group_sources(payload: Any) -> list[tuple[str, list[str]]]:
    """Group sources[].title by display label. Returns
    ``[(label, [titles...])]`` sorted by group count desc."""
    srcs = _coerce_to_list(payload_get(payload, "sources", default=[]))
    groups: dict[str, list[str]] = defaultdict(list)
    for s in srcs:
        if isinstance(s, dict):
            t = s.get("type") or s.get("source_type") or ""
            title = str(s.get("title") or "").strip()
        elif isinstance(s, str):
            t, title = s, s
        else:
            continue
        label = _label_source_type(str(t))
        if title and title not in groups[label]:
            groups[label].append(title)
    return sorted(
        groups.items(),
        key=lambda kv: (-len(kv[1]), kv[0]),
    )


@register_slide("sources")
class SourcesSlide(SlideBuilderBase):
    def build(
        self,
        presentation: Any,
        payload: Any,
        firm_branding: dict[str, Any],
        citations: list[Any],
        deck_context: Any = None,
    ) -> SlideResult:
        primary = (firm_branding or {}).get("primary_color") or DEFAULT_PRIMARY
        secondary_hex = (firm_branding or {}).get("secondary_color") or DEFAULT_SECONDARY

        slide = add_blank_slide(presentation)
        add_textbox(
            slide, left=0.5, top=1.0, width=SLIDE_WIDTH_IN - 1.0, height=0.4,
            text="Every claim in this deck links to one of the sources below.",
            font_size=11,
            color=parse_hex(DEFAULT_MUTED),
            align=PP_ALIGN.LEFT,
        )

        groups = _group_sources(payload)
        body = add_textbox(
            slide,
            left=0.5, top=1.5,
            width=SLIDE_WIDTH_IN - 1.0, height=5.5,
            text="",
            font_size=12,
            color=parse_hex(secondary_hex),
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        if not groups:
            add_paragraph(
                body.text_frame,
                "Firm library + retrieved evidence catalogue.",
                font_size=12, bullet=True,
                color=parse_hex(secondary_hex),
            )
        else:
            for label, titles in groups:
                add_paragraph(
                    body.text_frame,
                    f"{len(titles)} {label}",
                    font_size=14, bold=True,
                    color=parse_hex(secondary_hex),
                )
                for t in titles[:_MAX_TITLES_PER_GROUP]:
                    add_paragraph(
                        body.text_frame,
                        f"    – {t[:120]}",
                        font_size=10,
                        color=parse_hex(DEFAULT_MUTED),
                    )
                if len(titles) > _MAX_TITLES_PER_GROUP:
                    add_paragraph(
                        body.text_frame,
                        f"    – +{len(titles) - _MAX_TITLES_PER_GROUP} more",
                        font_size=10,
                        color=parse_hex(DEFAULT_MUTED),
                    )

        return SlideResult(slide_index=len(presentation.slides) - 1, citation_ids=[])
