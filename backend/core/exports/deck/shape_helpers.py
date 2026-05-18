"""Shape composition helpers for framework slides — W11/D3.

Centralizes the geometry/styling primitives the 2x2 + Porter's
slides share. Kept separate from ``_layout.py`` because these are
specifically shape-composition helpers (multi-shape constructs)
rather than single text-box wrappers.
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ._layout import (
    DEFAULT_FONT,
    DEFAULT_MUTED,
    DEFAULT_SECONDARY,
    parse_hex,
)

# Intensity → fill colour. Same palette as the W11/D2 risks_matrix
# severity column so the deck reads consistently.
_INTENSITY_COLOURS: dict[str, str] = {
    "low":      "#0F6E56",
    "moderate": "#B8860B",
    "medium":   "#B8860B",
    "high":     "#B91C1C",
}

# Minimum font size we'll render at before truncating. Per spec hard
# rule "Don't make text auto-shrink below 9pt."
_MIN_FONT_PT = 9


def add_quadrant_grid(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    x_axis_label: str,
    x_low_label: str,
    x_high_label: str,
    y_axis_label: str,
    y_low_label: str,
    y_high_label: str,
    primary_hex: str,
    muted_hex: str = DEFAULT_MUTED,
) -> dict[str, dict[str, float]]:
    """Draw a 2×2 grid + axis labels.

    Returns a dict keyed by quadrant name (``top_left``, ``top_right``,
    ``bottom_left``, ``bottom_right``) whose value is
    ``{left, top, width, height}`` in inches — the caller places items
    inside those regions.

    Geometry choices the spec authorizes us to make:
      - Y-axis label sits in the left margin (0.55 in wide), rotated
        90° via a separate textbox (pptx's text-rotation is brittle
        across renderers; we keep the label vertical-stacked instead).
      - Pole labels go just outside the grid on the corresponding edge.
      - The grid itself is 4 rectangles with thin borders rather than
        an inner cross, which renders more consistently across
        PowerPoint / Keynote / LibreOffice.
    """
    primary = parse_hex(primary_hex)
    muted = parse_hex(muted_hex)

    # Reserve margins for labels around the grid.
    label_pad_left = 0.7   # y-axis label column
    label_pad_right = 0.0
    label_pad_top = 0.45    # y high-pole label row
    label_pad_bottom = 0.55  # x-axis label row

    grid_left = left + label_pad_left
    grid_top = top + label_pad_top
    grid_width = width - label_pad_left - label_pad_right
    grid_height = height - label_pad_top - label_pad_bottom
    half_w = grid_width / 2
    half_h = grid_height / 2

    # Four quadrant rectangles. Borders only (transparent fill) so
    # item text inside reads without competing colour.
    for qx, qy in ((0, 0), (1, 0), (0, 1), (1, 1)):
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(grid_left + qx * half_w),
            Inches(grid_top + qy * half_h),
            Inches(half_w), Inches(half_h),
        )
        rect.fill.background()  # transparent
        rect.line.color.rgb = primary
        rect.line.width = Pt(0.75)

    # Y-axis label (left of grid, vertical-stacked because
    # cross-suite rotation is unreliable).
    y_axis_box = slide.shapes.add_textbox(
        Inches(left), Inches(grid_top),
        Inches(label_pad_left - 0.05), Inches(grid_height),
    )
    y_axis_box.text_frame.word_wrap = True
    y_axis_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = y_axis_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    # Render the label letter-by-letter on separate paragraphs so
    # it reads top-to-bottom without depending on rotated text.
    label_chars = list(y_axis_label.upper())
    for i, ch in enumerate(label_chars):
        para = p if i == 0 else y_axis_box.text_frame.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = ch
        run.font.name = DEFAULT_FONT
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = muted

    # Y high-pole label — above the top-left corner of the grid.
    _label_box(
        slide,
        left=grid_left, top=top, width=grid_width / 2, height=label_pad_top,
        text=y_high_label, color=muted, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Y low-pole label — below the bottom-left corner of the grid
    # (placed in the x-axis label band, left-aligned).
    _label_box(
        slide,
        left=grid_left, top=grid_top + grid_height,
        width=grid_width / 2, height=label_pad_bottom / 2,
        text=y_low_label, color=muted, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # X-axis main label + low/high poles below the grid.
    _label_box(
        slide,
        left=grid_left, top=grid_top + grid_height + 0.25,
        width=grid_width, height=label_pad_bottom - 0.25,
        text=x_axis_label.upper(),
        color=muted, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, bold=True,
    )
    _label_box(
        slide,
        left=grid_left, top=grid_top + grid_height,
        width=grid_width / 2, height=0.25,
        text=x_low_label, color=muted, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )
    _label_box(
        slide,
        left=grid_left + grid_width / 2, top=grid_top + grid_height,
        width=grid_width / 2, height=0.25,
        text=x_high_label, color=muted, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP,
    )

    # Quadrant regions in inches the caller can write into.
    return {
        "top_left":     {"left": grid_left,            "top": grid_top,           "width": half_w, "height": half_h},
        "top_right":    {"left": grid_left + half_w,   "top": grid_top,           "width": half_w, "height": half_h},
        "bottom_left":  {"left": grid_left,            "top": grid_top + half_h,  "width": half_w, "height": half_h},
        "bottom_right": {"left": grid_left + half_w,   "top": grid_top + half_h,  "width": half_w, "height": half_h},
    }


def add_intensity_badge(
    slide: Any,
    *,
    left: float, top: float,
    width: float = 0.7, height: float = 0.3,
    intensity: str,
) -> Any:
    """Coloured pill badge for low / moderate / high intensities.
    Returns the shape so the caller can group / position relative to it."""
    sev = (intensity or "").lower().strip()
    colour = _INTENSITY_COLOURS.get(sev, _INTENSITY_COLOURS["moderate"])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_hex(colour)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sev.upper() if sev else "—"
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = parse_hex("#FFFFFF")
    return shape


def add_citation_chip(
    slide: Any,
    *,
    left: float, top: float,
    number: int,
    claim_id: str = "",
    primary_hex: str = "#0F6E56",
    width: float = 0.32, height: float = 0.22,
) -> Any:
    """Small numbered superscript-style chip placed beside or below
    an item. Returns the shape so the caller can anchor / re-position.
    The chip carries the claim_id in its alt-text-equivalent via the
    title attribute (.name) so test inspection can map chip → id."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_hex(primary_hex)
    shape.line.fill.background()
    if claim_id:
        try:
            shape.name = f"chip-{claim_id[:60]}"
        except Exception:
            pass
    tf = shape.text_frame
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = parse_hex("#FFFFFF")
    return shape


def add_text_in_box(
    slide: Any,
    *,
    left: float, top: float, width: float, height: float,
    text: str,
    font_size: int = 11,
    bold: bool = False,
    color: RGBColor | None = None,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    font_name: str = DEFAULT_FONT,
    min_font_pt: int = _MIN_FONT_PT,
    char_capacity_at_min: int | None = None,
) -> Any:
    """Add a text box and shrink font to fit, never below ``min_font_pt``.
    If at the floor the text still overflows ``char_capacity_at_min``,
    truncate with an ellipsis rather than rendering unreadably small.
    """
    text = (text or "").strip()
    # Capacity heuristic: ~6 chars per inch at 9pt (Calibri). Caller
    # may override via char_capacity_at_min.
    if char_capacity_at_min is None:
        # rows visible at min font: height (inches) * 1.5 (lines/inch)
        chars_per_line = max(int(width * 9.0), 8)
        rows = max(int(height * 1.5), 1)
        char_capacity_at_min = chars_per_line * rows

    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor

    # Truncation pass (do this BEFORE shrinking — if it overflows at
    # the floor, the spec says truncate not shrink).
    chosen_font = font_size
    if len(text) > char_capacity_at_min:
        text = text[: char_capacity_at_min - 1].rstrip() + "…"
    # Naive shrink ladder: if requested font would clearly overflow,
    # step down until we hit the floor. We don't measure precisely
    # (pptx doesn't have a text-bounds API); the ladder is a defensive
    # downscale for long-text cases.
    soft_capacity = char_capacity_at_min * (chosen_font / min_font_pt)
    while chosen_font > min_font_pt and len(text) > soft_capacity:
        chosen_font -= 1
        soft_capacity = char_capacity_at_min * (chosen_font / min_font_pt)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(chosen_font)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = parse_hex(DEFAULT_SECONDARY)
    return shape


def _label_box(
    slide: Any,
    *,
    left: float, top: float, width: float, height: float,
    text: str, color: RGBColor,
    align: int, anchor: int,
    bold: bool = False, size: int = 10,
) -> Any:
    """Lightweight internal helper for axis pole labels — kept private
    because callers shouldn't need to drop one of these without going
    through ``add_quadrant_grid``."""
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = (text or "").strip()
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


__all__ = [
    "add_quadrant_grid",
    "add_intensity_badge",
    "add_citation_chip",
    "add_text_in_box",
]
