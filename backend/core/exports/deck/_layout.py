"""Shared layout + styling helpers for deck slides — W11/D1.

Centralises pptx unit math, color parsing, font/size constants, and
the colour-coding rules for recommendation verdicts. Per-mode slides
import from here to stay consistent without re-importing pptx's
verbose unit helpers each time.

Hard constraint per spec: no images embedded today, no external
asset fetching — logo URL is ignored on Day 1, falling back to
firm name as text on the title slide.
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Standard 16:9 slide is 13.333 x 7.5 inches.
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

# Default font stack — Calibri ships with PowerPoint/Office on macOS
# + Windows; LibreOffice falls back gracefully.
DEFAULT_FONT = "Calibri"

# Brand defaults if firms.branding is unset (matches the W10 1-pager
# defaults so HTML + PDF + PPTX visually agree).
DEFAULT_PRIMARY = "#0F6E56"
DEFAULT_SECONDARY = "#1B1F23"
DEFAULT_MUTED = "#5b6470"

# Recommendation verdict colour map. Mirrors W10's classify_recommendation
# rules so the deck's exec-summary panel matches the HTML/PDF 1-pager.
_VERDICT_COLOURS: dict[str, str] = {
    "green": "#0F6E56",
    "amber": "#B8860B",
    "red":   "#B91C1C",
    "neutral": "#1B1F23",
}


def parse_hex(value: Any, default: str = "#1B1F23") -> RGBColor:
    """Parse ``"#RRGGBB"`` (case-insensitive) into a pptx RGBColor.
    Falls back to ``default`` on bad input — never raises, because a
    bad branding row shouldn't crash the deck render."""
    if not isinstance(value, str):
        value = default
    s = value.strip().lstrip("#")
    if len(s) != 6:
        s = default.strip().lstrip("#")
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        d = default.strip().lstrip("#")
        return RGBColor(int(d[0:2], 16), int(d[2:4], 16), int(d[4:6], 16))


def verdict_colour(klass: str) -> RGBColor:
    """Map ``classify_recommendation`` outputs (green/amber/red/neutral)
    to a pptx RGBColor."""
    return parse_hex(_VERDICT_COLOURS.get(klass, DEFAULT_SECONDARY))


def set_slide_size_16_9(presentation: Any) -> None:
    """Force the presentation to 16:9 regardless of the default
    layout master, so slides render the same shape on PowerPoint,
    Keynote and LibreOffice."""
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)


def add_blank_slide(presentation: Any) -> Any:
    """Use the blank layout — index 6 by convention in the default
    pptx master. Falls back to the last available layout if a custom
    master has fewer slides."""
    layouts = presentation.slide_layouts
    blank_idx = 6 if len(layouts) > 6 else len(layouts) - 1
    return presentation.slides.add_slide(layouts[blank_idx])


def add_textbox(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor | None = None,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    font_name: str = DEFAULT_FONT,
    word_wrap: bool = True,
) -> Any:
    """Wrap pptx's verbose ``add_textbox`` + paragraph setup. Accepts
    inches as floats, returns the shape so the caller can mutate the
    text further (e.g. add multiple paragraphs)."""
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
    return shape


def add_paragraph(
    text_frame: Any,
    text: str,
    *,
    font_size: int = 12,
    bold: bool = False,
    color: RGBColor | None = None,
    bullet: bool = False,
    align: int = PP_ALIGN.LEFT,
    font_name: str = DEFAULT_FONT,
) -> Any:
    """Append a paragraph to an existing text frame. ``bullet`` is a
    textual indicator only (we prefix ``"• "``) — pptx's native
    bullet formatting is layout-master-dependent and inconsistent
    across renderers; a literal bullet character is portable."""
    p = text_frame.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = ("• " + text) if bullet else text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return p


# W11/D4: title bar reserved 0.65 in tall × full width. Footer reserved
# 0.32 in tall. Footnote strip sits just above the footer (max 0.45 in).
# Content sits between (top edge 0.85 in, bottom edge 6.85 in usable).
TITLE_BAR_HEIGHT_IN = 0.65
FOOTER_HEIGHT_IN = 0.32
FOOTNOTES_MAX_HEIGHT_IN = 0.45
CONTENT_TOP_IN = TITLE_BAR_HEIGHT_IN + 0.20  # 0.85
CONTENT_BOTTOM_IN = SLIDE_HEIGHT_IN - FOOTER_HEIGHT_IN - FOOTNOTES_MAX_HEIGHT_IN  # 6.73


def add_title_bar(
    slide: Any,
    *,
    title_text: str,
    primary_hex: str,
    secondary_hex: str = "#FFFFFF",
    font_size: int = 22,
) -> None:
    """Branded title bar — full width, height TITLE_BAR_HEIGHT_IN,
    primary-colour fill, white-on-primary bold text.

    Used by every content slide for visual coherence so the partner
    can flip through the deck and immediately know which firm it
    came from.
    """
    # Coloured band.
    band = add_horizontal_band(
        slide,
        left=0.0, top=0.0,
        width=SLIDE_WIDTH_IN, height=TITLE_BAR_HEIGHT_IN,
        color_hex=primary_hex,
    )
    # Title text inside the band — left-padded.
    shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.05),
        Inches(SLIDE_WIDTH_IN - 1.0), Inches(TITLE_BAR_HEIGHT_IN - 0.1),
    )
    tf = shape.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = (title_text or "").strip()
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = parse_hex(secondary_hex)


def add_footer(
    slide: Any,
    *,
    footer_text: str,
    page_number: int,
    total_pages: int,
    muted_hex: str = "#5b6470",
) -> None:
    """Footer strip — left text + right page number — at the bottom
    of every slide. Small grey 8pt monospace-ish styling so it reads
    "professional doc" rather than "PowerPoint default."""
    top = SLIDE_HEIGHT_IN - FOOTER_HEIGHT_IN
    # Footer text left.
    left_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(top),
        Inches(SLIDE_WIDTH_IN * 0.75), Inches(FOOTER_HEIGHT_IN),
    )
    tf = left_shape.text_frame
    tf.margin_left = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = (footer_text or "").strip()
    run.font.name = DEFAULT_FONT
    run.font.size = Pt(8)
    run.font.color.rgb = parse_hex(muted_hex)

    # Page number right.
    right_shape = slide.shapes.add_textbox(
        Inches(SLIDE_WIDTH_IN - 1.5), Inches(top),
        Inches(1.0), Inches(FOOTER_HEIGHT_IN),
    )
    tf2 = right_shape.text_frame
    tf2.margin_right = Inches(0.02)
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = f"{page_number} / {total_pages}"
    run2.font.name = DEFAULT_FONT
    run2.font.size = Pt(8)
    run2.font.color.rgb = parse_hex(muted_hex)


def add_citation_footnotes(
    slide: Any,
    *,
    footnotes: list[tuple[int, str]],
    muted_hex: str = "#5b6470",
) -> int:
    """Add a per-slide footnote strip just above the footer.

    ``footnotes`` is a list of ``(number, label)`` where number is the
    chip number rendered on the slide and label is the source-breadcrumb
    text (e.g. ``"SEC 10-K · Apple Inc. · Item 1A"``).

    Returns the count of footnotes actually rendered. If the joined
    string would exceed ~2 lines at 7pt, truncate with ``…`` rather
    than overflowing per spec hard rule.
    """
    if not footnotes:
        return 0
    parts = [f"^{n} {label}" for n, label in footnotes]
    joined = "    ".join(parts)
    # Heuristic character cap for ~2 lines at 7pt monospace across a
    # 13.333-in slide minus 1-in margins: ~280 chars/line × 2.
    max_chars = 560
    if len(joined) > max_chars:
        joined = joined[: max_chars - 1].rstrip() + "…"

    top = SLIDE_HEIGHT_IN - FOOTER_HEIGHT_IN - FOOTNOTES_MAX_HEIGHT_IN
    shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(top),
        Inches(SLIDE_WIDTH_IN - 1.0), Inches(FOOTNOTES_MAX_HEIGHT_IN),
    )
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = joined
    run.font.name = "Consolas, monospace"  # cosmetic, falls back if not present
    run.font.size = Pt(7)
    run.font.color.rgb = parse_hex(muted_hex)
    # Mark the shape so tests can find it by name.
    try:
        shape.name = "argus-citation-footnotes"
    except Exception:
        pass
    return len(footnotes)


def apply_theme_font(presentation: Any, font_name: str) -> None:
    """Set the slide-master's default font so every new text frame
    inherits the firm's font_family. python-pptx exposes
    ``slide_master.theme`` only partially — we walk the
    ``slide_master.element`` XML and set ``majorFont`` / ``minorFont``
    inside the ``fontScheme``.

    Per spec hard rule: don't replace the master theme entirely;
    only override fonts (and colour set further down). Falls back
    silently if the master XML doesn't expose the expected nodes.
    """
    try:
        master = presentation.slide_master
        # The slide-master XML uses the standard a: namespace.
        from pptx.oxml.ns import qn

        theme_elem = master.element.find(".//" + qn("a:theme"))
        if theme_elem is None:
            return
        font_scheme = theme_elem.find(".//" + qn("a:fontScheme"))
        if font_scheme is None:
            return
        for tag in ("a:majorFont", "a:minorFont"):
            section = font_scheme.find(qn(tag))
            if section is None:
                continue
            latin = section.find(qn("a:latin"))
            if latin is not None:
                latin.set("typeface", font_name)
    except Exception:
        # Theme XML shapes vary; never let a theme override break the
        # whole render. The per-shape font we set on add_textbox runs
        # is the fallback path that still works.
        pass


def add_horizontal_band(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    color_hex: str,
) -> Any:
    """Solid coloured rectangle, used for header bands + the
    recommendation verdict band on the exec-summary slide.

    We use a textbox with a solid fill rather than ``add_shape`` so
    we keep the same kind-of-element across slides (simplifies tests
    that inspect ``slide.shapes``)."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_hex(color_hex)
    shape.line.fill.background()
    return shape
