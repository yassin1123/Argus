"""DeckBuilder — wraps a pptx ``Presentation``, dispatches per-slide
builders by name, accumulates citations across slides, applies the
per-deck branding chrome (title bar / footer / footnotes), and
serializes the deck to bytes.

W11/D1: thin orchestrator + mode-specific sequencing.
W11/D2: 10 mode-specific slide builders plug into the registry.
W11/D3: framework visuals (2x2 + Porter's) on top of the same registry.
W11/D4: branded chrome — title bar with firm primary colour on every
        content slide, footer with firm footer_text + page number,
        per-slide citation footnotes mapped from a deck-wide registry.
        Slide builders opt out of the chrome via ``SlideResult.skip_chrome``
        (only the title slide does today).
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation

from ._layout import (
    DEFAULT_FONT,
    DEFAULT_PRIMARY,
    add_citation_footnotes,
    add_footer,
    add_title_bar,
    apply_theme_font,
    set_slide_size_16_9,
)
from .slides import get_slide_builder, list_registered_slides  # noqa: F401
from .slides._base import DeckContext, SlideResult


# Display titles for the title bar on each slide kind. Mode-specific
# slides override their own title text via the slide builder's
# add_title_bar call; this map is the fallback when a builder doesn't
# override.
_DEFAULT_TITLE_BAR_TEXT: dict[str, str] = {
    "title": "",  # title slide draws its own title; chrome is skipped
    "exec_summary": "Executive Summary",
    "context": "Context & Objectives",
    "target_overview": "Target Overview",
    "financial_profile": "Financial Profile",
    "valuation_range": "Valuation Range",
    "two_by_two_visual": "Strategic Options Matrix",
    "porters_five_forces_visual": "Industry Forces",
    "risks_matrix": "Risks & Mitigations",
    "integration_plan": "Integration Plan",
    "options_matrix": "Strategic Options",
    "market_landscape": "Market Landscape",
    "recommendation": "Recommendation",
    "next_steps": "Next Steps",
    "sources": "Sources",
}


class DeckBuilder:
    """Wraps a ``pptx.Presentation``. Each ``add_slide(name)`` call
    looks the builder up in the slide registry and runs it against
    the shared payload + branding + citations context.

    W11/D4 additions:
      - threads a :class:`DeckContext` into every slide builder for
        chip-number ↔ footnote synchronisation.
      - applies the branded chrome (title bar + footer + footnote
        strip) AFTER the slide builder runs, unless the builder
        returned ``SlideResult.skip_chrome=True``.
      - applies the firm's font family to the slide master once at
        init so every new text run inherits it.
    """

    def __init__(
        self,
        payload: Any,
        firm_branding: dict[str, Any] | None,
        citations: list[Any] | None,
    ) -> None:
        self._payload = payload
        self._branding = dict(firm_branding or {})
        self._citations = list(citations or [])
        self._presentation = Presentation()
        set_slide_size_16_9(self._presentation)
        # W11/D4: apply firm font_family to the slide master so newly
        # added text frames inherit it. Per-shape font names we set
        # in helpers stay authoritative on a per-shape basis.
        apply_theme_font(
            self._presentation,
            str(self._branding.get("font_family") or DEFAULT_FONT).split(",")[0].strip(),
        )

        self._results: list[SlideResult] = []
        self._slide_names_ordered: list[str] = []
        self._all_citation_ids: list[str] = []
        self._seen_citation_ids: set[str] = set()
        # Deck-wide context (citation registry + per-slide chip
        # numbers). Threaded into each slide builder.
        self._context = DeckContext()
        # Pre-populate breadcrumbs from any ClaimCitation objects the
        # service handed us, so chips on slides with no per-item
        # citations still get readable footnotes.
        for c in self._citations:
            cid = getattr(c, "claim_id", "") or ""
            if not cid:
                continue
            breadcrumb = self._breadcrumb_for_citation(c)
            if breadcrumb:
                self._context.citation_breadcrumbs[cid] = breadcrumb

    @staticmethod
    def _breadcrumb_for_citation(c: Any) -> str:
        """Human-readable source breadcrumb for a ClaimCitation."""
        source_type = (getattr(c, "source_type", "") or "").strip()
        title = (getattr(c, "source_title", "") or "").strip()
        parts: list[str] = []
        if source_type:
            parts.append(source_type.replace("_", " "))
        if title:
            parts.append(title)
        cid = (getattr(c, "claim_id", "") or "").strip()
        if cid:
            parts.append(cid)
        return " · ".join(parts) or cid

    @property
    def presentation(self) -> Presentation:
        return self._presentation

    @property
    def slide_count(self) -> int:
        return len(self._presentation.slides)

    @property
    def citation_count(self) -> int:
        return len(self._all_citation_ids)

    @property
    def cited_claim_ids(self) -> list[str]:
        return list(self._all_citation_ids)

    @property
    def slide_names(self) -> list[str]:
        """The ordered list of slide_name strings actually rendered."""
        return list(self._slide_names_ordered)

    @property
    def context(self) -> DeckContext:
        return self._context

    def add_slide(self, slide_name: str) -> SlideResult:
        builder_cls = get_slide_builder(slide_name)
        builder = builder_cls()

        # Open a fresh per-slide bucket on the context so chip
        # numbers register correctly.
        self._context.start_slide()

        result = builder.build(
            self._presentation,
            self._payload,
            self._branding,
            self._citations,
            self._context,
        )
        self._results.append(result)
        self._slide_names_ordered.append(slide_name)

        # Update deck-wide citation accumulator + ensure every cited
        # claim_id has a registered chip number even if the slide
        # builder didn't draw a shape chip (e.g. exec_summary uses
        # text-level superscripts). The chrome's footnote pass needs
        # a number-per-citation regardless of how it surfaced
        # visually.
        for cid in result.citation_ids:
            if not cid:
                continue
            if cid not in self._seen_citation_ids:
                self._seen_citation_ids.add(cid)
                self._all_citation_ids.append(cid)
            n = self._context.assign_chip(cid)
            self._context.record_chip_on_current_slide(n)

        return result

    def finalize_chrome(self) -> None:
        """Stamp title-bar, footer, and per-slide citation footnotes
        on every content slide once the sequence is fully assembled.

        Runs as a post-pass so the footer's ``page_number / total_pages``
        text knows the final slide count up front and so the
        slide-builder code itself doesn't need to take a dependency
        on the chrome helpers.
        """
        primary_hex = str(self._branding.get("primary_color") or DEFAULT_PRIMARY)
        footer_text = str(self._branding.get("footer_text") or "").strip()
        if not footer_text:
            firm = str(self._branding.get("_firm_name") or "Argus")
            footer_text = f"Confidential — Prepared by {firm}"

        total = self.slide_count
        for i, slide_result in enumerate(self._results, start=1):
            if slide_result.skip_chrome:
                continue
            slide = self._presentation.slides[slide_result.slide_index]
            slide_name = self._slide_names_ordered[slide_result.slide_index]
            # Title bar: pull explicit title text from the default map.
            title = _DEFAULT_TITLE_BAR_TEXT.get(slide_name, slide_name.replace("_", " ").title())
            if title:
                add_title_bar(slide, title_text=title, primary_hex=primary_hex)

            # Citation footnotes for this slide (resolved via context).
            footnotes = self._context_footnotes_for(slide_result.slide_index)
            if footnotes:
                add_citation_footnotes(slide, footnotes=footnotes)

            # Footer with firm text + page number.
            add_footer(
                slide,
                footer_text=footer_text,
                page_number=i,
                total_pages=total,
            )

    def _context_footnotes_for(self, slide_index: int) -> list[tuple[int, str]]:
        if slide_index >= len(self._context.per_slide_chip_numbers):
            return []
        chips = self._context.per_slide_chip_numbers[slide_index]
        rev = {n: cid for cid, n in self._context.citation_numbers.items()}
        out: list[tuple[int, str]] = []
        for n in sorted(set(chips)):
            cid = rev.get(n, "")
            label = self._context.citation_breadcrumbs.get(cid) or cid or "—"
            out.append((n, label))
        return out

    def serialize(self) -> bytes:
        # W11/D4: always finalize chrome before serializing. Builders
        # may call this explicitly (tests), but the serialize path is
        # the contract surface.
        self.finalize_chrome()
        buf = io.BytesIO()
        self._presentation.save(buf)
        return buf.getvalue()
