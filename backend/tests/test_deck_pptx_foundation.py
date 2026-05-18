"""Phase 3 / Week 11 / Day 1 — deck PPTX foundation tests.

Six tests per spec covering the round-trip, content presence on
each base slide, and the slide-count contract.

All tests run as pure unit tests against
``DeckPptxExporter.render(...)`` — no DB, no FastAPI, no IO beyond
optional disk writes for the openable-via-reopen test.
"""

from __future__ import annotations

import io
import tempfile
from pathlib import Path
from typing import Any

import pytest
from pptx import Presentation

from core.exports._base import ClaimCitation
from core.exports.deck_pptx import DeckPptxExporter
from core.exports.deck.sequences import get_deck_sequence_for_mode


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


_BRANDING: dict[str, Any] = {
    "primary_color": "#0F6E56",
    "secondary_color": "#1B1F23",
    "footer_text": "Test Firm · Confidential",
    "_firm_name": "Test Firm",
}


def _m_and_a_payload() -> dict[str, Any]:
    return {
        "mode": "m_and_a_diligence",
        "recommendation": "PROCEED WITH CONDITIONS",
        "confidence_level": "Medium-High",
        "key_reasons": [
            "Stable cash flow from anchor customers.",
            "Segment leadership in facilities maintenance.",
            "Synergy potential of £6.5m run-rate.",
            "Brand premium attached to historic franchise.",
        ],
        "risks": [
            "Customer concentration in top-3 = 41%.",
            "Halo contract renewal binary in March 2026.",
            "ROI segment EBITDA-negative.",
        ],
        "sources": [
            {"type": "firm_library", "title": "M&A Target Screen Playbook"},
            {"type": "sec_filing", "title": "10-K 2023"},
            {"type": "earnings_transcript", "title": "Q4 2024 call"},
            {"type": "firm_library", "title": "Retail Sector Primer"},
        ],
        "valuation_range": {
            "low": {"gbp_m": 205.0}, "base": {"gbp_m": 220.0}, "high": {"gbp_m": 235.0},
        },
        "deal_structure_implications": {
            "walk_away_triggers": [
                "If 'Project Halo' contract is not renewed, walk.",
                "If FY25 EBITDA tracks under £18m at H1.",
            ],
        },
        "recommendation_claim_ids": ["claim_1", "claim_2", "claim_3"],
        "_engagement_title": "TargetCo Holdings M&A diligence",
        "_target_name": "TargetCo Holdings",
        "_firm_name": "Test Firm",
    }


def _citations(n: int) -> list[ClaimCitation]:
    return [
        ClaimCitation(
            claim_id=f"claim_{i+1}",
            text=f"Claim text #{i+1} grounded in evidence.",
            source_title=f"Source doc {i+1}",
            source_type="firm_library",
        )
        for i in range(n)
    ]


def _all_text(slide: Any) -> str:
    """Concatenate every run on every text frame on the slide."""
    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text:
                    parts.append(run.text)
    return "\n".join(parts)


@pytest.fixture
def exporter() -> DeckPptxExporter:
    return DeckPptxExporter()


# ---------------------------------------------------------------------------
# Test 1 — round trip: bytes → reopen → resolved-sequence slide count
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_renders_pptx_round_trip(exporter: DeckPptxExporter) -> None:
    """Verify the bytes round-trip back through python-pptx. W11/D2
    expanded the M&A sequence from 3 to 10 slides; this test now
    checks against the live ``get_deck_sequence_for_mode`` result so
    the assertion tracks future sequence additions automatically."""
    result = await exporter.render(_m_and_a_payload(), _BRANDING, _citations(6))
    assert result.file_size > 0
    # PPTX is a ZIP — header magic is PK\x03\x04.
    assert result.file_bytes[:4] == b"PK\x03\x04"
    prs = Presentation(io.BytesIO(result.file_bytes))
    expected_sequence = get_deck_sequence_for_mode("m_and_a_diligence")
    assert len(prs.slides) == len(expected_sequence)
    assert result.metadata["mode"] == "m_and_a_diligence"
    assert result.metadata["slide_sequence"] == expected_sequence
    # Spot-check the three base-trio slides land in the sequence.
    for base in ("title", "exec_summary", "recommendation"):
        assert base in expected_sequence


# ---------------------------------------------------------------------------
# Test 2 — title slide contains the recommendation text
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_title_slide_contains_recommendation_text(
    exporter: DeckPptxExporter,
) -> None:
    result = await exporter.render(_m_and_a_payload(), _BRANDING, _citations(3))
    prs = Presentation(io.BytesIO(result.file_bytes))
    title_text = _all_text(prs.slides[0])
    assert "PROCEED WITH CONDITIONS" in title_text
    # Firm name surfaces in the prepared-by line.
    assert "Test Firm" in title_text


# ---------------------------------------------------------------------------
# Test 3 — exec summary slide has three columns (with headings + body)
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_exec_summary_three_columns(exporter: DeckPptxExporter) -> None:
    result = await exporter.render(_m_and_a_payload(), _BRANDING, _citations(5))
    prs = Presentation(io.BytesIO(result.file_bytes))
    text = _all_text(prs.slides[1])
    # Heading + three column titles must all be present.
    assert "Executive Summary" in text
    assert "Recommendation" in text
    assert "Top reasons" in text
    assert "Top risks" in text
    # And a sample of body content from each column.
    assert "PROCEED" in text  # recommendation column
    assert "Stable cash flow" in text  # reasons column
    assert "Customer concentration" in text  # risks column


# ---------------------------------------------------------------------------
# Test 4 — recommendation slide has full prose + source-panel content
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_recommendation_slide_has_panel_text(
    exporter: DeckPptxExporter,
) -> None:
    """Recommendation slide moved index when W11/D2 expanded the
    sequence. Look it up by name via ``slide_sequence`` rather than
    a hardcoded index."""
    result = await exporter.render(_m_and_a_payload(), _BRANDING, _citations(3))
    prs = Presentation(io.BytesIO(result.file_bytes))
    seq = result.metadata["slide_sequence"]
    idx = seq.index("recommendation")
    text = _all_text(prs.slides[idx])
    assert "Recommendation" in text
    # Full recommendation prose lands on this slide.
    assert "PROCEED WITH CONDITIONS" in text
    # M&A walk-away trigger surfaces on the recommendation slide.
    assert "Walk-away triggers" in text
    assert "Project Halo" in text
    # Source panel aggregates by source label.
    assert "Sources" in text
    assert "SEC filings" in text or "firm-library" in text


# ---------------------------------------------------------------------------
# Test 5 — slide count matches the resolved sequence
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_slide_count_matches_sequence(exporter: DeckPptxExporter) -> None:
    """Slide count == resolved-sequence length, for every supported
    mode. W11/D2 added mode-specific sequences (M&A=10, growth=9,
    general=7) on top of the W11/D1 base trio (title, exec_summary,
    recommendation), and unknown modes still fall back to general."""
    payload = _m_and_a_payload()
    result = await exporter.render(payload, _BRANDING, _citations(2))
    prs = Presentation(io.BytesIO(result.file_bytes))
    expected = get_deck_sequence_for_mode("m_and_a_diligence")
    assert len(prs.slides) == len(expected)
    assert result.metadata["slide_count"] == len(expected)
    # The base trio is invariant across every mode.
    for mode in ("m_and_a_diligence", "growth_strategy", "general"):
        seq = get_deck_sequence_for_mode(mode)
        for base in ("title", "exec_summary", "recommendation"):
            assert base in seq, f"{mode} missing base slide {base}"
    # Unknown mode falls back to general.
    assert (
        get_deck_sequence_for_mode("does_not_exist")
        == get_deck_sequence_for_mode("general")
    )


# ---------------------------------------------------------------------------
# Test 6 — file writes to disk and reopens without error
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_pptx_opens_without_error(
    exporter: DeckPptxExporter, tmp_path: Path
) -> None:
    """Write the bytes to a real .pptx path and reopen via
    ``Presentation(filepath)``. python-pptx's loader does its own
    XML validation pass, so this catches anything that round-trips
    through BytesIO but fails on a fresh open."""
    result = await exporter.render(_m_and_a_payload(), _BRANDING, _citations(4))
    fpath = tmp_path / "test_deck.pptx"
    fpath.write_bytes(result.file_bytes)
    assert fpath.exists()
    assert fpath.stat().st_size == result.file_size
    # Reopen — raises if the file is malformed.
    prs = Presentation(str(fpath))
    assert len(prs.slides) == len(get_deck_sequence_for_mode("m_and_a_diligence"))
    # File size sanity. W11/D1 capped at 100KB on the 3-slide deck;
    # W11/D2 expanded to 10 slides + native column chart + tables —
    # 200KB is the new headroom for the same no-images contract.
    # Real-world payloads (W7 demo) come in around 50KB even with
    # 10 slides; the cap is a sanity check, not a tuning target.
    assert result.file_size < 200_000, (
        f"deck file too large: {result.file_size} bytes (limit 200KB after W11/D2)"
    )
