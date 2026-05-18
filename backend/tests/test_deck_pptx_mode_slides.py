"""Phase 3 / Week 11 / Day 2 — mode-specific slide tests.

Nine tests per spec covering:
  - mode-aware sequence dispatch (M&A vs growth_strategy vs general)
  - mode-specific slide content (target overview, valuation, integration,
    market landscape)
  - mode exclusion (growth deck must not include valuation; M&A deck
    must not include market_landscape)
  - chart-data round-trip on the financial_profile slide
  - severity colour-coding on the risks_matrix table

Tests don't depend on the DB — they exercise ``DeckPptxExporter`` +
``DeckBuilder`` directly with hand-rolled payloads, and re-parse the
resulting pptx to assert content.
"""

from __future__ import annotations

import io
from typing import Any

import pytest
from pptx import Presentation

from core.exports._base import ClaimCitation
from core.exports.deck.sequences import get_deck_sequence_for_mode
from core.exports.deck_pptx import DeckPptxExporter


_BRANDING: dict[str, Any] = {
    "primary_color": "#0F6E56",
    "secondary_color": "#1B1F23",
    "_firm_name": "Test Firm",
}


def _m_and_a_payload() -> dict[str, Any]:
    return {
        "mode": "m_and_a_diligence",
        "recommendation": "PROCEED WITH CONDITIONS",
        "confidence_level": "Medium-High",
        "summary": "TargetCo stable; deal de-risks via earnouts.",
        "key_reasons": ["Stable cash flow.", "Segment leadership.", "Synergy 6.5m run-rate."],
        "risks": ["Concentration."],
        "sources": [
            {"type": "firm_library", "title": "M&A Playbook"},
            {"type": "firm_library", "title": "TargetCo CIM"},
            {"type": "sec_filing", "title": "10-K 2023"},
            {"type": "earnings_transcript", "title": "Q4 2024 call"},
        ],
        "target_overview": {
            "name": "TargetCo",
            "business_model": "UK industrial services provider with facilities + mechanical segments.",
            "segments": [
                {"name": "Facilities", "revenue_pct": 52.0, "growth_rate": "+2.8%"},
                {"name": "Mechanical", "revenue_pct": 28.0, "growth_rate": "+1.2%"},
            ],
            "geographies": [{"geography": "UK", "revenue_pct": 91.0}, {"geography": "ROI", "revenue_pct": 9.0}],
            "ownership_history": "Founder-owned then PE-backed.",
            "key_customers_concentration": "Top 3 = 41% of revenue.",
        },
        "financial_profile": {
            "revenue_trajectory": {"points": [
                {"period": "FY21", "value_gbp_m": 153.2, "source_citation": "x"},
                {"period": "FY22", "value_gbp_m": 168.5, "source_citation": "x"},
                {"period": "FY23", "value_gbp_m": 190.0, "source_citation": "x"},
                {"period": "FY24", "value_gbp_m": 203.0, "source_citation": "x"},
            ]},
            "ebitda_trajectory": {"points": [
                {"period": "FY23", "value_gbp_m": 19.0, "source_citation": "x"},
                {"period": "FY24", "value_gbp_m": 21.5, "source_citation": "x"},
            ]},
            "margin_profile": {"gross_margin": "36.4%", "ebitda_margin": "10.6%", "fcf_margin": "6.2%"},
            "working_capital_dynamics": "51-day cycle.",
            "debt_structure": "Net debt 32m.",
            "capex_intensity": "4.5% of revenue.",
            "cash_flow_quality": "Recurring.",
        },
        "valuation_range": {
            "low":  {"gbp_m": 205.0, "methodology": "DCF @ WACC 10%", "key_assumptions": ["a1", "a2"]},
            "base": {"gbp_m": 220.0, "methodology": "EV/EBITDA 8.5x", "key_assumptions": ["b1", "b2"]},
            "high": {"gbp_m": 235.0, "methodology": "EV/Sales 1.4x", "key_assumptions": ["c1", "c2"]},
            "comparable_transactions_cited": [
                {"target": "Comp A", "acquirer": "PE A", "year": 2023, "multiple": "8.2x", "source_citation": "db"},
            ],
        },
        "risks_and_mitigations": [
            {"risk_category": "commercial", "description": "Halo renewal binary.", "severity": "high", "mitigation": "Earnout.", "residual_risk": "Timing"},
            {"risk_category": "operational", "description": "ROI EBITDA-negative.", "severity": "medium", "mitigation": "Restructure.", "residual_risk": "12-mo payback"},
            {"risk_category": "financial", "description": "Customer concentration.", "severity": "low", "mitigation": "Diversify.", "residual_risk": "2yr horizon"},
        ],
        "integration_plan": {
            "day_one_priorities": ["Customer comms", "Payroll continuity"],
            "first_100_days": [
                {"workstream": "IT", "owner_role": "CTO", "milestone": "Email migration"},
                {"workstream": "GTM", "owner_role": "CRO", "milestone": "Cross-sell map"},
            ],
            "first_year": [{"workstream": "Operations", "owner_role": "COO", "milestone": "Plant consolidation"}],
            "integration_complexity_rating": "medium",
            "complexity_rationale": "moderate overlap",
        },
        "next_steps": [{"action": "Confirm Halo renewal.", "owner_role": "Lead Partner", "timing": "2 weeks"}],
        "recommendation_claim_ids": ["claim_1", "claim_2"],
        "_engagement_title": "TargetCo M&A diligence",
        "_target_name": "TargetCo Holdings",
        "_firm_name": "Test Firm",
    }


def _growth_payload() -> dict[str, Any]:
    return {
        "mode": "growth_strategy",
        "recommendation": "Launch Scotland pilot before North-East entry.",
        "confidence_level": "Medium",
        "summary": "Scotland de-risks expansion.",
        "key_reasons": ["Existing base", "Channel proven", "Lower capex"],
        "risks": ["Halo binary", "Capital trade-off"],
        "sources": [{"type": "firm_library", "title": "Growth Framework"}, {"type": "document", "title": "Capex"}],
        "evidence_ledger_summary": "Strong UK channel-access evidence; weaker on TAM size.",
        "counterarguments": ["Aldi/Lidl discounter pressure", "Online entrants"],
        "next_steps": [{"action": "Sign 3 anchor customers by month 6.", "owner_role": "CRO", "timing": "Q1 2026"}],
        "options_matrix": [
            {"name": "Scotland pilot", "rationale": "Existing base", "quadrant": "top_right"},
            {"name": "North-East entry", "rationale": "Higher TAM", "quadrant": "bottom_right"},
        ],
        "_engagement_title": "TargetCo Scotland pilot",
        "_firm_name": "Test Firm",
    }


def _slide_names_from_result(result: Any) -> list[str]:
    """Read the slide sequence the DeckPptxExporter records on the
    ExporterResult metadata. (pptx Slide objects reject arbitrary
    attribute assignment, so the builder tracks names separately.)"""
    seq = result.metadata.get("slide_sequence") or []
    assert isinstance(seq, list), "slide_sequence must be a list"
    return list(seq)


def _all_text(slide: Any) -> str:
    """Concatenate every run on every text frame on the slide,
    including text inside table cells."""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        parts.append(run.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                parts.append(run.text)
    return "\n".join(parts)


def _all_text_in_tables(slide: Any) -> list[list[str]]:
    """Extract a 2D matrix of cell texts from every table on the slide."""
    out: list[list[str]] = []
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        cell_text += r.text
                row_texts.append(cell_text.strip())
            out.append(row_texts)
    return out


@pytest.fixture
def exporter() -> DeckPptxExporter:
    return DeckPptxExporter()


# ---------------------------------------------------------------------------
# Test 1 — M&A deck includes target_overview at the expected index
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_m_and_a_deck_has_target_overview_slide(exporter: DeckPptxExporter) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    # Mode shape (live-checked rather than hardcoded so future sequence
    # additions don't break this test). M&A must include all of:
    # title, exec_summary, target_overview, financial_profile,
    # valuation_range, risks_matrix, integration_plan, recommendation,
    # next_steps, sources — plus optionally framework visuals (W11/D3).
    required = (
        "title", "exec_summary", "target_overview", "financial_profile",
        "valuation_range", "risks_matrix", "integration_plan",
        "recommendation", "next_steps", "sources",
    )
    for s in required:
        assert s in names, f"M&A sequence missing required slide {s!r}"
    # And it doesn't sprout growth-only slides.
    for s in ("market_landscape", "context", "options_matrix"):
        assert s not in names, f"M&A sequence wrongly includes growth-only slide {s!r}"
    idx = names.index("target_overview")
    text = _all_text(prs.slides[idx])
    assert "Target Overview" in text
    # Business model + segment + ownership content all surface.
    assert "industrial services" in text.lower()
    assert "Facilities" in text  # segment table content
    assert "ownership" in text.lower() or "founder-owned" in text.lower()


# ---------------------------------------------------------------------------
# Test 2 — valuation_range slide has all three boxes with their numbers
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_m_and_a_deck_has_valuation_range_with_three_boxes(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    idx = names.index("valuation_range")
    text = _all_text(prs.slides[idx])
    # Three scenario headings present (LOW / BASE / HIGH).
    assert "LOW" in text
    assert "BASE" in text
    assert "HIGH" in text
    # The three monetary values land verbatim (writer payload: 205 / 220 / 235).
    assert "£205.0m" in text
    assert "£220.0m" in text
    assert "£235.0m" in text
    # Methodology and a comparable-transactions line surface.
    assert "DCF" in text or "WACC" in text
    assert "comparable" in text.lower() or "8.2x" in text


# ---------------------------------------------------------------------------
# Test 3 — integration_plan slide has three bands
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_m_and_a_deck_has_integration_plan_three_bands(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    idx = names.index("integration_plan")
    text = _all_text(prs.slides[idx])
    for band in ("Day 1 priorities", "First 100 days", "First year"):
        assert band in text, f"missing band: {band}"
    # Real content from each band surfaces.
    assert "Payroll continuity" in text  # day-1
    assert "Email migration" in text     # first 100 days
    assert "Plant consolidation" in text  # first year
    # Complexity badge present.
    assert "COMPLEXITY" in text
    assert "MEDIUM" in text


# ---------------------------------------------------------------------------
# Test 4 — growth deck includes market_landscape slide
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_growth_deck_has_market_landscape(exporter: DeckPptxExporter) -> None:
    r = await exporter.render(_growth_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    # Live-checked: growth deck must include all required structural
    # slides. options_matrix was replaced by porters_five_forces_visual
    # in W11/D3, so accept either as the strategic-options slot.
    required = (
        "title", "exec_summary", "context", "market_landscape",
        "recommendation", "risks_matrix", "next_steps", "sources",
    )
    for s in required:
        assert s in names, f"growth sequence missing required slide {s!r}"
    # Strategic-options slot is either the W11/D3 visual or the old D2 stub.
    assert "porters_five_forces_visual" in names or "options_matrix" in names
    idx = names.index("market_landscape")
    text = _all_text(prs.slides[idx])
    assert "Market Landscape" in text
    # Either real narrative or the documented fallback string lands;
    # in this fixture we provided evidence_ledger_summary so narrative
    # should win.
    assert "Strong UK channel-access" in text or "Market overview" in text


# ---------------------------------------------------------------------------
# Test 5 — growth deck omits valuation slide (mode exclusivity)
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_growth_deck_omits_valuation_slide(exporter: DeckPptxExporter) -> None:
    r = await exporter.render(_growth_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    assert "valuation_range" not in names
    # And the M&A-only target_overview / financial_profile / integration_plan
    # likewise stay out.
    assert "target_overview" not in names
    assert "financial_profile" not in names
    assert "integration_plan" not in names


# ---------------------------------------------------------------------------
# Test 6 — M&A deck omits market_landscape (mode exclusivity, other side)
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_m_and_a_deck_omits_market_landscape(exporter: DeckPptxExporter) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    assert "market_landscape" not in names
    assert "options_matrix" not in names
    assert "context" not in names  # M&A uses target_overview instead


# ---------------------------------------------------------------------------
# Test 7 — general mode falls back to the minimal 7-slide sequence
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_general_mode_falls_back_to_minimal_sequence(
    exporter: DeckPptxExporter,
) -> None:
    payload = {
        "mode": "general",
        "recommendation": "Pursue option A.",
        "key_reasons": ["r1", "r2"],
        "risks": ["x1"],
        "sources": [],
        "next_steps": [{"action": "Do thing", "owner_role": "Lead", "timing": "Q2"}],
        "_engagement_title": "Generic engagement",
    }
    r = await exporter.render(payload, _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    assert names == get_deck_sequence_for_mode("general")
    assert len(names) == 7
    # And M&A / growth-specific slides aren't present.
    for excluded in (
        "target_overview", "financial_profile", "valuation_range",
        "integration_plan", "market_landscape", "options_matrix",
    ):
        assert excluded not in names


# ---------------------------------------------------------------------------
# Test 8 — financial_profile chart carries the trajectory's data points
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_financial_profile_chart_has_data_points(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    slide = prs.slides[names.index("financial_profile")]

    # Find the chart shape and pull out the underlying series data.
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break
    assert chart_shape is not None, "financial_profile slide must have a chart"

    chart = chart_shape.chart
    # Categories are the periods.
    cats = list(chart.plots[0].categories)
    assert cats == ["FY21", "FY22", "FY23", "FY24"]
    # Series values match payload.
    series_values = list(chart.plots[0].series[0].values)
    assert series_values == pytest.approx([153.2, 168.5, 190.0, 203.0])


# ---------------------------------------------------------------------------
# Test 9 — risks_matrix renders severity per row with colour-coded cells
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_risks_matrix_renders_severity_color_per_row(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    names = _slide_names_from_result(r)
    slide = prs.slides[names.index("risks_matrix")]

    rows = _all_text_in_tables(slide)
    # ``_all_text_in_tables`` returns a flat list of rows across every
    # table on the slide. The risks_matrix slide has exactly one table
    # (header + 3 body rows).
    assert rows, "risks_matrix slide must contain at least one table"
    assert rows[0] == ["Risk", "Severity", "Mitigation"]
    severities = [r[1] for r in rows[1:]]
    assert severities == ["HIGH", "MEDIUM", "LOW"]
    # Spot-check the colour fill on the severity cells.
    risk_table = None
    for shape in slide.shapes:
        if shape.has_table:
            risk_table = shape.table
            break
    assert risk_table is not None
    expected_fills = {
        "HIGH":   (0xB9, 0x1C, 0x1C),
        "MEDIUM": (0xB8, 0x86, 0x0B),
        "LOW":    (0x0F, 0x6E, 0x56),
    }
    for i, (sev, rgb) in enumerate(zip(severities, [expected_fills[s] for s in severities]), start=1):
        cell = risk_table.cell(i, 1)
        actual = cell.fill.fore_color.rgb
        assert (actual[0], actual[1], actual[2]) == rgb, (
            f"row {i} severity {sev}: expected fill {rgb}, got {(actual[0], actual[1], actual[2])}"
        )
