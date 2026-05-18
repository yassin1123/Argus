"""Phase 3 / Week 11 / Day 3 — framework visual tests.

Eight tests per spec covering:
  - 2x2 grid + quadrant placement + axis labels
  - Porter's diagram + intensity badge colours
  - Both visuals' fallback behaviour when the framework is absent
  - Citation chips render on every framework slide
"""

from __future__ import annotations

import io
from typing import Any

import pytest
from pptx import Presentation
from pptx.util import Emu

from core.exports.deck.shape_helpers import _INTENSITY_COLOURS
from core.exports.deck_pptx import DeckPptxExporter


_BRANDING: dict[str, Any] = {
    "primary_color": "#0F6E56",
    "secondary_color": "#1B1F23",
    "_firm_name": "Test Firm",
}


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _m_and_a_payload(*, with_two_by_two: bool = True) -> dict[str, Any]:
    base: dict[str, Any] = {
        "mode": "m_and_a_diligence",
        "recommendation": "PROCEED WITH CONDITIONS",
        "confidence_level": "Medium-High",
        "summary": "TargetCo stable.",
        "key_reasons": ["r1", "r2", "r3"],
        "risks": ["x1"],
        "sources": [{"type": "firm_library", "title": "Playbook"}],
        "target_overview": {
            "name": "TargetCo",
            "business_model": "UK industrial services.",
            "segments": [{"name": "Facilities", "revenue_pct": 52.0, "growth_rate": "+2.8%"}],
            "geographies": [{"geography": "UK", "revenue_pct": 91.0}],
            "ownership_history": "Founder-owned.",
            "key_customers_concentration": "Top 3 = 41%.",
        },
        "financial_profile": {
            "revenue_trajectory": {"points": [
                {"period": "FY23", "value_gbp_m": 190.0, "source_citation": "x"},
                {"period": "FY24", "value_gbp_m": 203.0, "source_citation": "x"},
            ]},
            "ebitda_trajectory": {"points": [
                {"period": "FY23", "value_gbp_m": 19.0, "source_citation": "x"},
                {"period": "FY24", "value_gbp_m": 21.5, "source_citation": "x"},
            ]},
            "margin_profile": {"gross_margin": "36.4%", "ebitda_margin": "10.6%", "fcf_margin": "6.2%"},
            "working_capital_dynamics": "51-day cycle.",
            "debt_structure": "Net debt 32m.",
            "capex_intensity": "4.5%.",
            "cash_flow_quality": "Recurring.",
        },
        "valuation_range": {
            "low":  {"gbp_m": 205.0, "methodology": "DCF", "key_assumptions": ["a"]},
            "base": {"gbp_m": 220.0, "methodology": "EV/EBITDA", "key_assumptions": ["b"]},
            "high": {"gbp_m": 235.0, "methodology": "EV/Sales", "key_assumptions": ["c"]},
        },
        "risks_and_mitigations": [
            {"risk_category": "commercial", "description": "Halo binary.",
             "severity": "high", "mitigation": "Earnout.", "residual_risk": "Timing"},
        ],
        "integration_plan": {
            "day_one_priorities": ["Customer comms"],
            "first_100_days": [{"workstream": "IT", "owner_role": "CTO", "milestone": "Email"}],
            "first_year": [{"workstream": "Ops", "owner_role": "COO", "milestone": "Consolidate"}],
            "integration_complexity_rating": "medium",
            "complexity_rationale": "moderate",
        },
        "next_steps": [{"action": "Confirm Halo renewal.", "owner_role": "Lead Partner", "timing": "2 weeks"}],
        "_engagement_title": "TargetCo M&A diligence",
        "_target_name": "TargetCo Holdings",
        "_firm_name": "Test Firm",
    }
    if with_two_by_two:
        base["frameworks"] = {
            "two_by_two": {
                "title": "TargetCo deal screen",
                "x_axis_label": "Strategic fit", "x_axis_low_label": "Low", "x_axis_high_label": "High",
                "y_axis_label": "Deal complexity", "y_axis_low_label": "Low", "y_axis_high_label": "High",
                "items": [
                    {"name": "Facilities maintenance", "quadrant": "top_right",
                     "rationale": "High fit + moderate complexity.",
                     "evidence_citations": ["claim_a1"]},
                    {"name": "ROI segment", "quadrant": "bottom_right",
                     "rationale": "Moderate fit, high complexity.",
                     "evidence_citations": ["claim_a2"]},
                    {"name": "Carve-out", "quadrant": "top_left",
                     "rationale": "Lower fit, simpler integration.",
                     "evidence_citations": ["claim_a3"]},
                    {"name": "Walk away", "quadrant": "bottom_left",
                     "rationale": "Default if Halo not renewed.",
                     "evidence_citations": ["claim_a4"]},
                ],
                "interpretation": "Facilities is the deal driver; ROI is contingent.",
            },
        }
    return base


def _growth_payload(*, with_porters: bool = True) -> dict[str, Any]:
    base: dict[str, Any] = {
        "mode": "growth_strategy",
        "recommendation": "Launch Scotland pilot Q1.",
        "confidence_level": "Medium",
        "summary": "Scotland de-risks expansion.",
        "key_reasons": ["Existing base", "Channel proven", "Lower capex"],
        "risks": ["Halo binary"],
        "sources": [{"type": "firm_library", "title": "Growth Framework"}],
        "next_steps": [{"action": "Sign 3 anchor customers.", "owner_role": "CRO", "timing": "Q1"}],
        "_engagement_title": "TargetCo Scotland pilot",
        "_firm_name": "Test Firm",
    }
    if with_porters:
        base["frameworks"] = {
            "porters_five_forces": {
                "market_definition": "UK industrial services.",
                "rivalry": {"intensity": "high", "rationale": "Established players.",
                            "key_drivers": ["Market saturation", "Price competition"],
                            "evidence_citations": ["claim_p_riv"]},
                "supplier_power": {"intensity": "moderate", "rationale": "Specialised suppliers.",
                                    "key_drivers": ["Concentration"],
                                    "evidence_citations": ["claim_p_sup"]},
                "buyer_power": {"intensity": "high", "rationale": "Large clients.",
                                "key_drivers": ["Concentration"],
                                "evidence_citations": ["claim_p_buy"]},
                "substitute_threat": {"intensity": "low", "rationale": "Few substitutes.",
                                       "key_drivers": ["Service uniqueness"],
                                       "evidence_citations": ["claim_p_sub"]},
                "new_entrant_threat": {"intensity": "moderate", "rationale": "Some barriers.",
                                        "key_drivers": ["Capital"],
                                        "evidence_citations": ["claim_p_ent"]},
                "overall_attractiveness": "moderate",
                "overall_rationale": "Rivalry + buyer power offset by low substitute threat.",
            },
        }
    return base


@pytest.fixture
def exporter() -> DeckPptxExporter:
    return DeckPptxExporter()


# ---------------------------------------------------------------------------
# Helpers — find slides by name, walk shapes, harvest rectangles + chips
# ---------------------------------------------------------------------------


def _slide_index(result: Any, name: str) -> int:
    seq = result.metadata.get("slide_sequence") or []
    return seq.index(name)


def _all_text(slide: Any) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        parts.append(run.text)
    return "\n".join(parts)


def _rectangles(slide: Any) -> list[Any]:
    """Return shapes whose auto-shape type is rectangle. Lets us count
    the four quadrant rectangles without picking up title bands etc."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out = []
    for s in slide.shapes:
        try:
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.auto_shape_type is not None:
                out.append(s)
        except Exception:
            pass
    return out


def _chip_shapes(slide: Any) -> list[Any]:
    """Citation chips are named 'chip-<claim_id>' by ``add_citation_chip``."""
    return [s for s in slide.shapes if getattr(s, "name", "").startswith("chip-")]


# ---------------------------------------------------------------------------
# Test 1 — 2x2 visual renders four quadrant rectangles
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_two_by_two_visual_renders_four_quadrants(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    slide = prs.slides[_slide_index(r, "two_by_two_visual")]
    from pptx.enum.shapes import MSO_SHAPE
    rects = [s for s in _rectangles(slide) if s.auto_shape_type == MSO_SHAPE.RECTANGLE]
    # Top brand band + 4 quadrant rectangles = at least 5 rectangles.
    # Strict count would be brittle (other helpers may add bands too);
    # check that at least 4 distinct quadrant rectangles are present.
    assert len(rects) >= 5, f"expected ≥5 rectangles (1 band + 4 quadrants), got {len(rects)}"
    # The four quadrants should be roughly equal in size and lie inside
    # the body region (top > 1.0 in, height < 3.0 in).
    quad_candidates = [r_ for r_ in rects if r_.top > Emu(1_000_000) and r_.height < Emu(3_000_000)]
    assert len(quad_candidates) >= 4, "expected at least 4 quadrant-shaped rectangles in body region"


# ---------------------------------------------------------------------------
# Test 2 — items land in their correct quadrants
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_two_by_two_visual_renders_items_in_correct_quadrants(
    exporter: DeckPptxExporter,
) -> None:
    payload = _m_and_a_payload()
    r = await exporter.render(payload, _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    slide = prs.slides[_slide_index(r, "two_by_two_visual")]

    # Pull the centre of the 4-cell grid. ``add_quadrant_grid`` places
    # the grid at left=0.4+0.7=1.1in, top=1.2+0.45=1.65in, etc. To
    # avoid coupling tests to exact px, we infer the grid centre from
    # the rectangles themselves.
    from pptx.enum.shapes import MSO_SHAPE
    rects = [
        s for s in _rectangles(slide)
        if s.auto_shape_type == MSO_SHAPE.RECTANGLE and s.height < Emu(3_000_000)
        and s.top > Emu(800_000)
    ]
    # 4 quadrant rects (the band is shorter; we filtered to body region).
    assert len(rects) >= 4
    rects_sorted = sorted(rects, key=lambda s: (s.top, s.left))[:4]
    # Mid-x / mid-y is the average of the rectangle centres.
    centres = [(r_.left + r_.width // 2, r_.top + r_.height // 2) for r_ in rects_sorted]
    xs = sorted(set(c[0] for c in centres))
    ys = sorted(set(c[1] for c in centres))
    grid_mid_x = (xs[0] + xs[-1]) // 2
    grid_mid_y = (ys[0] + ys[-1]) // 2

    # For each payload item, find a text shape whose run contains the
    # item name and verify its centre is in the matching quadrant.
    items = payload["frameworks"]["two_by_two"]["items"]
    name_to_quadrant = {it["name"]: it["quadrant"] for it in items}

    def _shape_text(s: Any) -> str:
        if not s.has_text_frame:
            return ""
        return "".join(r.text for p in s.text_frame.paragraphs for r in p.runs)

    for name, expected_quadrant in name_to_quadrant.items():
        match = None
        for s in slide.shapes:
            txt = _shape_text(s)
            if name in txt and len(txt) < 200:  # the item-name shape, not the interpretation
                match = s
                break
        assert match is not None, f"item {name!r} not found on slide"
        cx = match.left + match.width // 2
        cy = match.top + match.height // 2
        is_top = cy < grid_mid_y
        is_left = cx < grid_mid_x
        actual = ("top" if is_top else "bottom") + "_" + ("left" if is_left else "right")
        assert actual == expected_quadrant, (
            f"item {name!r} expected in {expected_quadrant}, found in {actual} "
            f"(centre {cx},{cy} vs grid mid {grid_mid_x},{grid_mid_y})"
        )


# ---------------------------------------------------------------------------
# Test 3 — 2x2 visual includes axis labels
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_two_by_two_visual_includes_axes_labels(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    slide = prs.slides[_slide_index(r, "two_by_two_visual")]
    text = _all_text(slide)
    # X-axis main label rendered upper-cased.
    assert "STRATEGIC FIT" in text.upper()
    # The four pole labels.
    assert "Low" in text and "High" in text
    # Y-axis label rendered letter-stacked — the joined letters
    # should reproduce the original word when concatenated.
    # Look for the sequence DEAL... in the rendered text (vertical-stacked).
    y_label_letters = "".join(ch for ch in text.upper() if ch.isalpha())
    assert "DEALCOMPLEXITY" in y_label_letters


# ---------------------------------------------------------------------------
# Test 4 — Porter's diagram renders five force boxes
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_porters_visual_renders_five_force_boxes(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_growth_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    slide = prs.slides[_slide_index(r, "porters_five_forces_visual")]
    text = _all_text(slide)
    for label in (
        "Competitive Rivalry",
        "Threat of New Entrants",
        "Threat of Substitutes",
        "Supplier Power",
        "Buyer Power",
    ):
        assert label in text, f"missing Porter's label {label!r}"
    # Overall attractiveness footer.
    assert "Overall attractiveness" in text or "OVERALL" in text.upper()


# ---------------------------------------------------------------------------
# Test 5 — Porter's intensity badges colour-match the intensity level
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_porters_intensity_badges_color_correct(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_growth_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    slide = prs.slides[_slide_index(r, "porters_five_forces_visual")]

    # Find every ROUNDED_RECTANGLE shape — the intensity-badge primitive.
    from pptx.enum.shapes import MSO_SHAPE
    badges = []
    for s in slide.shapes:
        try:
            if s.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                badges.append(s)
        except Exception:
            pass

    # 5 force badges + 1 overall-attractiveness badge = 6.
    assert len(badges) >= 5, f"expected ≥5 intensity badges, got {len(badges)}"

    # For each badge, read its text + fill rgb, then assert the fill
    # matches the canonical palette.
    expected_rgb = {
        "low":      (0x0F, 0x6E, 0x56),
        "moderate": (0xB8, 0x86, 0x0B),
        "medium":   (0xB8, 0x86, 0x0B),
        "high":     (0xB9, 0x1C, 0x1C),
    }
    for badge in badges:
        text = ""
        if badge.has_text_frame:
            for p in badge.text_frame.paragraphs:
                for r_ in p.runs:
                    text += r_.text
        text_norm = text.strip().lower()
        if text_norm not in expected_rgb:
            continue
        fill = badge.fill.fore_color.rgb
        assert (fill[0], fill[1], fill[2]) == expected_rgb[text_norm], (
            f"badge {text_norm!r} expected {expected_rgb[text_norm]}, got {(fill[0], fill[1], fill[2])}"
        )


# ---------------------------------------------------------------------------
# Test 6 — 2x2 visual falls back when payload lacks the framework
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_two_by_two_falls_back_when_payload_lacks_framework(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_m_and_a_payload(with_two_by_two=False), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    slide = prs.slides[_slide_index(r, "two_by_two_visual")]
    text = _all_text(slide)
    assert "Strategic Options Matrix" in text
    assert "not produced for this engagement" in text.lower()


# ---------------------------------------------------------------------------
# Test 7 — Porter's visual falls back when payload lacks the framework
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_porters_falls_back_when_payload_lacks_framework(
    exporter: DeckPptxExporter,
) -> None:
    r = await exporter.render(_growth_payload(with_porters=False), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    slide = prs.slides[_slide_index(r, "porters_five_forces_visual")]
    text = _all_text(slide)
    assert "Industry Forces" in text or "Porter" in text
    assert "not produced for this engagement" in text.lower()


# ---------------------------------------------------------------------------
# Test 8 — both framework slides carry citation chips
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_each_framework_slide_has_citation_chips(
    exporter: DeckPptxExporter,
) -> None:
    # M&A — 4 items each with one citation -> 4 chips on the 2x2 slide.
    r_mna = await exporter.render(_m_and_a_payload(), _BRANDING, [])
    prs_mna = Presentation(io.BytesIO(r_mna.file_bytes))
    two_by_two_slide = prs_mna.slides[_slide_index(r_mna, "two_by_two_visual")]
    chips = _chip_shapes(two_by_two_slide)
    assert len(chips) >= 4, (
        f"expected ≥4 citation chips on 2x2 slide, got {len(chips)}"
    )
    # The chip names embed the claim_id (chip-<claim_id>); verify our
    # 4 payload claim_ids are represented.
    chip_names = {s.name for s in chips}
    for cid in ("claim_a1", "claim_a2", "claim_a3", "claim_a4"):
        assert any(cid in name for name in chip_names), f"no chip for {cid!r}"

    # Growth — 5 force boxes each with one citation chip.
    r_g = await exporter.render(_growth_payload(), _BRANDING, [])
    prs_g = Presentation(io.BytesIO(r_g.file_bytes))
    porters_slide = prs_g.slides[_slide_index(r_g, "porters_five_forces_visual")]
    p_chips = _chip_shapes(porters_slide)
    assert len(p_chips) >= 5, f"expected ≥5 citation chips on Porter's slide, got {len(p_chips)}"


# ---------------------------------------------------------------------------
# Bonus sanity: registry exposes the canonical intensity palette.
# ---------------------------------------------------------------------------


def test_intensity_palette_is_consistent() -> None:
    """Spec hard rule: deck intensity colours match the risks-matrix
    palette (low=green, moderate=amber, high=red). This is a guardrail
    test so future contributors don't drift the palette in one place."""
    assert _INTENSITY_COLOURS["low"] == "#0F6E56"
    assert _INTENSITY_COLOURS["high"] == "#B91C1C"
    assert _INTENSITY_COLOURS["moderate"] == "#B8860B"
