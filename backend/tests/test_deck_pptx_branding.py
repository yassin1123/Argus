"""Phase 3 / Week 11 / Day 4 — branding + citation footnotes tests.

Eight tests per spec covering:
  - asset cache: fetch once + 24h TTL + Pillow resize cap
  - title slide branding (firm primary colour)
  - title bar on every content slide
  - footer on every slide
  - per-slide citation footnotes mirror chip count
  - logo fetch failure falls back to firm-name text
  - footnote truncation when too many citations
"""

from __future__ import annotations

import asyncio
import io
import re
from pathlib import Path
from typing import Any
from unittest import mock
from uuid import uuid4

import pytest
from PIL import Image
from pptx import Presentation

from core.exports import asset_cache
from core.exports._base import ClaimCitation
from core.exports.deck._layout import (
    FOOTER_HEIGHT_IN,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TITLE_BAR_HEIGHT_IN,
    add_citation_footnotes,
    parse_hex,
)
from core.exports.deck.deck_builder import DeckBuilder
from core.exports.deck_pptx import DeckPptxExporter


_BRANDING: dict[str, Any] = {
    "primary_color": "#0F6E56",
    "secondary_color": "#1B1F23",
    "font_family": "Calibri",
    "footer_text": "Test Firm · Confidential",
    "_firm_name": "Test Firm",
}


def _m_and_a_payload() -> dict[str, Any]:
    return {
        "mode": "m_and_a_diligence",
        "recommendation": "PROCEED WITH CONDITIONS",
        "confidence_level": "Medium-High",
        "summary": "TargetCo stable.",
        "key_reasons": ["Stable cash flow.", "Segment leadership.", "Synergy potential."],
        "risks": ["Customer concentration."],
        "sources": [{"type": "firm_library", "title": "M&A Playbook"}],
        "target_overview": {
            "name": "TargetCo", "business_model": "Industrial services.",
            "segments": [{"name": "Facilities", "revenue_pct": 52.0, "growth_rate": "+2.8%"}],
            "geographies": [{"geography": "UK", "revenue_pct": 91.0}],
            "ownership_history": "Founder-owned.", "key_customers_concentration": "Top 3 = 41%.",
        },
        "financial_profile": {
            "revenue_trajectory": {"points": [
                {"period": "FY23", "value_gbp_m": 190.0, "source_citation": "x"},
                {"period": "FY24", "value_gbp_m": 203.0, "source_citation": "x"},
            ]},
            "ebitda_trajectory": {"points": [
                {"period": "FY23", "value_gbp_m": 19.0, "source_citation": "x"},
                {"period": "FY24", "value_gbp_m": 21.5, "source_citation": "x"},
            ]},
            "margin_profile": {"gross_margin": "36.4%", "ebitda_margin": "10.6%", "fcf_margin": "6.2%"},
            "working_capital_dynamics": "51-day.", "debt_structure": "Net debt 32m.",
            "capex_intensity": "4.5%.", "cash_flow_quality": "Recurring.",
        },
        "valuation_range": {
            "low":  {"gbp_m": 205.0, "methodology": "DCF", "key_assumptions": ["a"]},
            "base": {"gbp_m": 220.0, "methodology": "EV/EBITDA", "key_assumptions": ["b"]},
            "high": {"gbp_m": 235.0, "methodology": "EV/Sales", "key_assumptions": ["c"]},
        },
        "risks_and_mitigations": [
            {"risk_category": "commercial", "description": "Halo renewal.",
             "severity": "high", "mitigation": "Earnout.", "residual_risk": "Timing"},
        ],
        "integration_plan": {
            "day_one_priorities": ["Customer comms"],
            "first_100_days": [{"workstream": "IT", "owner_role": "CTO", "milestone": "Email"}],
            "first_year": [{"workstream": "Ops", "owner_role": "COO", "milestone": "Consolidate"}],
            "integration_complexity_rating": "medium", "complexity_rationale": "moderate",
        },
        "next_steps": [{"action": "Confirm Halo renewal.", "owner_role": "Lead Partner", "timing": "2 weeks"}],
        "frameworks": {
            "two_by_two": {
                "title": "Test 2x2", "x_axis_label": "X", "x_axis_low_label": "Low", "x_axis_high_label": "High",
                "y_axis_label": "Y", "y_axis_low_label": "Low", "y_axis_high_label": "High",
                "items": [
                    {"name": "Item A", "quadrant": "top_right", "rationale": "RA",
                     "evidence_citations": ["claim_a1"]},
                    {"name": "Item B", "quadrant": "bottom_right", "rationale": "RB",
                     "evidence_citations": ["claim_a2"]},
                    {"name": "Item C", "quadrant": "top_left", "rationale": "RC",
                     "evidence_citations": ["claim_a3"]},
                    {"name": "Item D", "quadrant": "bottom_left", "rationale": "RD",
                     "evidence_citations": ["claim_a4"]},
                ],
                "interpretation": "Items in their cells.",
            },
        },
        "recommendation_claim_ids": ["claim_1"],
        "_engagement_title": "TargetCo M&A diligence",
        "_target_name": "TargetCo",
        "_firm_name": "Test Firm",
    }


# ---------------------------------------------------------------------------
# Test 1 — logo cached after first fetch (no second HTTP call)
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_logo_cached_after_first_fetch(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path,
) -> None:
    # Build a 50x50 red PNG so Pillow has something real to resize.
    raw_img = io.BytesIO()
    Image.new("RGBA", (50, 50), color=(255, 0, 0, 255)).save(raw_img, format="PNG")
    raw_bytes = raw_img.getvalue()

    monkeypatch.setattr(asset_cache, "LOGO_CACHE_DIR", tmp_path)
    fetcher = mock.AsyncMock(return_value=raw_bytes)
    monkeypatch.setattr(asset_cache, "_fetch_url", fetcher)

    firm_id = uuid4()
    url = "https://example.com/logo.png"
    out1 = await asset_cache.fetch_and_cache_logo(firm_id, url)
    out2 = await asset_cache.fetch_and_cache_logo(firm_id, url)
    assert out1 is not None and out2 is not None
    # Identical bytes on the second call (cache hit).
    assert out1 == out2
    # And the HTTP fetcher only ran once.
    assert fetcher.call_count == 1


# ---------------------------------------------------------------------------
# Test 2 — large source image gets resized to max 300px wide
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_logo_resized_to_max_width(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path,
) -> None:
    big = io.BytesIO()
    Image.new("RGBA", (1200, 600), color=(0, 110, 86, 255)).save(big, format="PNG")
    monkeypatch.setattr(asset_cache, "LOGO_CACHE_DIR", tmp_path)
    monkeypatch.setattr(asset_cache, "_fetch_url", mock.AsyncMock(return_value=big.getvalue()))

    out = await asset_cache.fetch_and_cache_logo(uuid4(), "https://example.com/big.png")
    assert out is not None
    im = Image.open(io.BytesIO(out))
    assert im.width <= 300, f"expected resize to ≤300px wide, got {im.width}"
    # Aspect ratio preserved (originally 2:1).
    assert abs((im.width / im.height) - 2.0) < 0.05, (
        f"expected ~2:1 aspect ratio, got {im.width}:{im.height}"
    )


# ---------------------------------------------------------------------------
# Test 3 — title slide uses firm primary colour on the recommendation text
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_title_slide_uses_firm_primary_color() -> None:
    r = await DeckPptxExporter().render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    title_slide = prs.slides[0]

    # Find the recommendation textbox by content + colour. The
    # title-slide recommendation is rendered at font_size 36, bold.
    expected_rgb = (0x0F, 0x6E, 0x56)
    matched = False
    for shape in title_slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "PROCEED" not in (run.text or ""):
                    continue
                fg = run.font.color
                try:
                    rgb = fg.rgb
                except Exception:
                    continue
                if (rgb[0], rgb[1], rgb[2]) == expected_rgb:
                    matched = True
    assert matched, "title-slide recommendation text not rendered in firm primary colour"


# ---------------------------------------------------------------------------
# Test 4 — title bar with primary-coloured background on every content slide
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_title_bar_on_content_slides() -> None:
    r = await DeckPptxExporter().render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    expected_rgb = (0x0F, 0x6E, 0x56)
    n_content_slides = 0
    n_with_title_bar = 0
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue  # title slide opts out (skip_chrome=True)
        n_content_slides += 1
        # Look for a shape at top=0 with primary-colour fill spanning
        # the full slide width.
        for shape in slide.shapes:
            if shape.top != 0:
                continue
            try:
                fill_rgb = shape.fill.fore_color.rgb
            except Exception:
                continue
            if (fill_rgb[0], fill_rgb[1], fill_rgb[2]) == expected_rgb:
                # Width should span (close to) the full slide.
                from pptx.util import Inches
                slide_width_emu = Inches(SLIDE_WIDTH_IN)
                if shape.width >= slide_width_emu * 0.95:
                    n_with_title_bar += 1
                    break
    assert n_content_slides > 0
    assert n_with_title_bar == n_content_slides, (
        f"only {n_with_title_bar}/{n_content_slides} content slides have a title bar"
    )


# ---------------------------------------------------------------------------
# Test 5 — footer (firm text + page number) on every slide
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_footer_on_every_slide() -> None:
    r = await DeckPptxExporter().render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    total = len(prs.slides)
    pagenum_pat = re.compile(r"^\d+\s*/\s*\d+$")
    for i, slide in enumerate(prs.slides):
        if i == 0:
            # Title slide carries its own "Confidential. Prepared by ..."
            # line; the page-number footer is content-slide only.
            continue
        texts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = ""
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    txt += run.text
            texts.append(txt.strip())
        assert any("Test Firm" in t and "Confidential" in t for t in texts), (
            f"slide {i} missing firm footer text"
        )
        assert any(pagenum_pat.match(t) for t in texts), (
            f"slide {i} missing page-number footer; texts={texts}"
        )
        # The expected page number for this slide.
        expected = f"{i + 1} / {total}"
        assert any(t == expected for t in texts), (
            f"slide {i} page number incorrect (expected {expected!r}, got {texts})"
        )


# ---------------------------------------------------------------------------
# Test 6 — citation footnote count matches chip count on a chip-bearing slide
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_citation_footnotes_match_chip_count() -> None:
    r = await DeckPptxExporter().render(_m_and_a_payload(), _BRANDING, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    seq = r.metadata["slide_sequence"]
    two_by_two_idx = seq.index("two_by_two_visual")
    slide = prs.slides[two_by_two_idx]

    # Count chip shapes (named ``chip-{claim_id}``).
    chips = [s for s in slide.shapes if getattr(s, "name", "").startswith("chip-")]
    assert chips, "expected chips on the two_by_two_visual slide"

    # Footnotes shape is the named ``argus-citation-footnotes`` textbox.
    foot_shape = next(
        (s for s in slide.shapes if getattr(s, "name", "") == "argus-citation-footnotes"),
        None,
    )
    assert foot_shape is not None, "missing argus-citation-footnotes shape on slide"

    foot_text = ""
    for p in foot_shape.text_frame.paragraphs:
        for run in p.runs:
            foot_text += run.text
    # Each footnote begins with the ``^N`` chip marker the helper writes.
    n_footnotes = foot_text.count("^")
    assert n_footnotes == len(chips), (
        f"chip count {len(chips)} != footnote count {n_footnotes} "
        f"(footnote text: {foot_text!r})"
    )


# ---------------------------------------------------------------------------
# Test 7 — logo fetch failure falls back to firm-name text
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_logo_fetch_failure_falls_back_to_text(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path,
) -> None:
    # Empty cache + fetch returns None (HTTP failure / 404).
    monkeypatch.setattr(asset_cache, "LOGO_CACHE_DIR", tmp_path)
    monkeypatch.setattr(asset_cache, "_fetch_url", mock.AsyncMock(return_value=None))

    branding = dict(_BRANDING)
    branding["logo_url"] = "https://nope.example.com/missing.png"
    branding["_firm_id"] = uuid4()

    r = await DeckPptxExporter().render(_m_and_a_payload(), branding, [])
    prs = Presentation(io.BytesIO(r.file_bytes))
    title_slide = prs.slides[0]
    texts: list[str] = []
    has_picture = False
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in title_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_picture = True
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if run.text:
                        texts.append(run.text)
    assert not has_picture, "expected no picture when fetch fails"
    joined = "\n".join(texts)
    assert "Test Firm" in joined


# ---------------------------------------------------------------------------
# Test 8 — citation footnote text truncates when too many citations
# ---------------------------------------------------------------------------


def test_footnote_truncated_when_too_many_citations(tmp_path: Path) -> None:
    """Unit-test the add_citation_footnotes helper directly: 20+
    long-breadcrumb citations exceed the 2-line cap (~560 chars) and
    must be truncated with an ellipsis.
    """
    from pptx import Presentation as _P
    from pptx.util import Inches as _I

    prs = _P()
    prs.slide_width = _I(SLIDE_WIDTH_IN)
    prs.slide_height = _I(SLIDE_HEIGHT_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    footnotes = [
        (i, f"Source Type · Long-title-document-name-{i} · Subsection-{i}")
        for i in range(1, 21)  # 20 entries
    ]
    n = add_citation_footnotes(slide, footnotes=footnotes)
    assert n == 20  # all entries logically requested

    # Find the rendered shape and check the text was truncated with "…".
    foot_shape = next(
        (s for s in slide.shapes if getattr(s, "name", "") == "argus-citation-footnotes"),
        None,
    )
    assert foot_shape is not None
    rendered = ""
    for p in foot_shape.text_frame.paragraphs:
        for run in p.runs:
            rendered += run.text
    assert rendered.endswith("…"), (
        f"expected ellipsis at end of truncated footnotes, got: {rendered[-40:]!r}"
    )
    # Truncated to roughly the helper's 560-char cap (allow a few extra
    # for the trailing ellipsis).
    assert len(rendered) <= 580, f"rendered footnote too long: {len(rendered)} chars"
