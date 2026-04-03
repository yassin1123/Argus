"""Template-driven PowerPoint: blueprint schema, rules, tables, fixed layout."""

import io
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field, field_validator

from deliverables.assemble import build_deliverable_document
from deliverables.models import DeliverableDocument


class SlideTableSpec(BaseModel):
    headers: list[str] = Field(default_factory=list)
    rows: list[list[str]] = Field(default_factory=list)

    model_config = {"extra": "ignore"}


class SlideSpec(BaseModel):
    """Title = conclusion-style headline; max 5 bullets; no paragraphs."""

    title: str = ""
    bullets: list[str] = Field(default_factory=list)
    table: SlideTableSpec | None = None
    layout_kind: str = Field(
        default="insight",
        description="title_cover | insight | matrix_table | evidence_appendix | risks",
    )

    model_config = {"extra": "ignore"}

    @field_validator("bullets", mode="before")
    @classmethod
    def clamp_bullets(cls, v: Any) -> list[str]:
        if v is None:
            return []
        if not isinstance(v, list):
            return []
        out: list[str] = []
        for x in v[:5]:
            s = str(x).strip().replace("\n\n", " ").replace("\n", " ")
            if len(s) > 320:
                s = s[:317] + "…"
            if s:
                out.append(s)
        return out


class SlideBlueprint(BaseModel):
    slides: list[SlideSpec] = Field(default_factory=list)
    version: int = 2

    model_config = {"extra": "ignore"}


def _truncate_title(s: str, n: int = 80) -> str:
    s = str(s).strip().replace("\n", " ")
    return (s[: n - 1] + "…") if len(s) > n else s


def build_slide_blueprint_from_document(
    *,
    doc: DeliverableDocument,
    report: dict[str, Any],
    session_query: str,
) -> SlideBlueprint:
    slides: list[SlideSpec] = []

    slides.append(
        SlideSpec(
            layout_kind="title_cover",
            title=_truncate_title(doc.cover_title, 72),
            bullets=[doc.cover_subtitle[:200] if doc.cover_subtitle else session_query[:200]],
        )
    )

    exec_bullets = list(doc.exec_insights)[:5]
    if not exec_bullets and report.get("summary"):
        exec_bullets = [str(report.get("summary", ""))[:400]]
    slides.append(
        SlideSpec(
            layout_kind="insight",
            title="Executive summary — key takeaways",
            bullets=exec_bullets or ["See detailed findings in appendix."],
        )
    )

    slides.append(
        SlideSpec(
            layout_kind="insight",
            title=_truncate_title(doc.exec_recommendation, 80) or "Recommendation",
            bullets=[str(x)[:300] for x in (report.get("key_reasons") or [])[:4] if str(x).strip()],
        )
    )

    for f in doc.findings[:8]:
        bl = [f.explanation[:280]] if f.explanation else []
        if f.mini_conclusion and f.mini_conclusion not in (bl[0] if bl else ""):
            bl.append(f.mini_conclusion[:200])
        slides.append(
            SlideSpec(
                layout_kind="insight",
                title=_truncate_title(f.title, 72),
                bullets=bl[:5],
            )
        )

    cp = report.get("consulting_payload") if isinstance(report.get("consulting_payload"), dict) else {}
    dc = cp.get("decision_criteria") if isinstance(cp.get("decision_criteria"), list) else []
    if dc:
        headers = ["Criterion", "Weight", "How met"]
        rows: list[list[str]] = []
        for row in dc[:8]:
            if not isinstance(row, dict):
                continue
            rows.append(
                [
                    str(row.get("criterion", ""))[:80],
                    str(row.get("weight", ""))[:40],
                    str(row.get("how_met", ""))[:120],
                ]
            )
        if rows:
            slides.append(
                SlideSpec(
                    layout_kind="matrix_table",
                    title="Decision criteria — structured view",
                    bullets=[],
                    table=SlideTableSpec(headers=headers, rows=rows),
                )
            )

    om = cp.get("options_matrix") if isinstance(cp.get("options_matrix"), list) else []
    if om:
        headers = ["Option", "Fit", "Pros", "Cons"]
        rows = []
        for row in om[:6]:
            if not isinstance(row, dict):
                continue
            pros = ", ".join(str(x) for x in (row.get("pros") or [])[:3])[:100]
            cons = ", ".join(str(x) for x in (row.get("cons") or [])[:3])[:100]
            rows.append(
                [
                    str(row.get("option", ""))[:60],
                    str(row.get("fit", ""))[:80],
                    pros,
                    cons,
                ]
            )
        if rows:
            slides.append(
                SlideSpec(
                    layout_kind="matrix_table",
                    title="Options comparison matrix",
                    bullets=[],
                    table=SlideTableSpec(headers=headers, rows=rows),
                )
            )

    appendix_lines = doc.appendix_sources[:12]
    if appendix_lines:
        slides.append(
            SlideSpec(
                layout_kind="evidence_appendix",
                title="Sources & evidence appendix",
                bullets=[s[:280] for s in appendix_lines],
            )
        )

    risk_bullets = [str(x)[:300] for x in doc.risks_body[:5] if str(x).strip()]
    if risk_bullets:
        slides.append(
            SlideSpec(layout_kind="risks", title="Key risks to monitor", bullets=risk_bullets)
        )

    return SlideBlueprint(slides=slides)


def build_slide_blueprint(
    *,
    report: dict[str, Any],
    session_query: str,
    session_title: str,
) -> SlideBlueprint:
    doc = build_deliverable_document(
        report=report, session_query=session_query, session_title=session_title
    )
    return build_slide_blueprint_from_document(
        doc=doc, report=report, session_query=session_query
    )


def _load_brand() -> dict[str, Any]:
    p = Path(__file__).resolve().parent / "brand_tokens.json"
    if not p.is_file():
        return {"accent_rgb": [99, 102, 241], "text_primary": "#0f172a"}
    return json.loads(p.read_text(encoding="utf-8"))


def _accent_rgb(brand: dict[str, Any]) -> RGBColor:
    rgb = brand.get("accent_rgb") or [99, 102, 241]
    return RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def render_pptx_from_blueprint(bp: SlideBlueprint) -> tuple[bytes, dict[str, Any]]:
    prs = Presentation()
    slides_meta: list[dict[str, Any]] = []
    brand = _load_brand()
    accent = _accent_rgb(brand)
    text_dark = RGBColor(15, 23, 42)

    for spec in bp.slides:
        if spec.table and spec.table.rows:
            layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(layout)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            tf.text = spec.title or "Slide"
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = text_dark

            nrows = len(spec.table.rows) + 1
            ncols = len(spec.table.headers)
            if ncols == 0:
                ncols = 1
            left, top, width, height = Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)
            table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
            for j, h in enumerate(spec.table.headers[:ncols]):
                cell = table.cell(0, j)
                cell.text = str(h)[:200]
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent
            for i, row in enumerate(spec.table.rows[:15]):
                for j in range(ncols):
                    cell = table.cell(i + 1, j)
                    val = str(row[j])[:300] if j < len(row) else ""
                    cell.text = val
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = text_dark
            slides_meta.append({"title": spec.title, "table": True, "layout": spec.layout_kind})
            continue

        if spec.layout_kind == "title_cover":
            layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(layout)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.8), Inches(10), Inches(0.08)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.9), Inches(8.5), Inches(1.2))
            tbf = tb.text_frame
            tbf.text = spec.title or "Argus"
            tbf.paragraphs[0].font.size = Pt(36)
            tbf.paragraphs[0].font.bold = True
            tbf.paragraphs[0].font.color.rgb = text_dark

            sub = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(8.5), Inches(2))
            sf = sub.text_frame
            bullets = spec.bullets[:3]
            sf.text = bullets[0] if bullets else " "
            sf.paragraphs[0].font.size = Pt(16)
            sf.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)
            for line in bullets[1:]:
                p = sf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(100, 116, 139)

            foot = slide.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(8.5), Inches(0.4))
            ff = foot.text_frame
            ff.text = "Argus · confidential"
            ff.paragraphs[0].font.size = Pt(10)
            ff.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)
            slides_meta.append({"title": spec.title, "layout": "title_cover"})
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = spec.title or "Slide"
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.color.rgb = text_dark
            if spec.layout_kind == "risks":
                p.font.color.rgb = RGBColor(185, 28, 28)
        body = slide.shapes.placeholders[1].text_frame
        bullets = spec.bullets[:5]
        if bullets:
            body.text = bullets[0]
            for line in bullets[1:]:
                p = body.add_paragraph()
                p.text = line
                p.level = 0
                p.font.size = Pt(14)
                p.font.color.rgb = text_dark
        else:
            body.text = " "
        slides_meta.append({"title": spec.title, "bullets": bullets, "layout": spec.layout_kind})

    bio = io.BytesIO()
    prs.save(bio)
    meta = {"slides": slides_meta, "version": bp.version, "brand": brand.get("name", "Argus")}
    return bio.getvalue(), meta


def build_and_render_pptx(
    *,
    report: dict[str, Any],
    session_query: str,
    session_title: str,
) -> tuple[bytes, dict[str, Any]]:
    bp = build_slide_blueprint(
        report=report, session_query=session_query, session_title=session_title
    )
    return render_pptx_from_blueprint(bp)
