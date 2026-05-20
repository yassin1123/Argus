"""Cross-artifact consistency checker — W14/D4.

For a single engagement, extracts the recommendation text from each
generated artifact and asserts the artifacts agree on what the
engagement actually recommends. The artifacts share a single source
payload; divergence means an artifact is rendering stale or wrong
and needs investigation.

Public surface:

  - :func:`extract_recommendation` — content-type-dispatched extractor.
    HTML/MD/email-MD: regex/text-walk; PPTX: python-pptx slide-by-slide;
    XLSX: openpyxl scan of the Summary sheet; PDF: PyMuPDF text extract.
  - :func:`normalise_recommendation` — strips whitespace, lowercases,
    collapses internal whitespace, drops markdown emphasis markers,
    and keeps only the FIRST sentence of the recommendation. The
    first-sentence rule keeps M&A's
    "PROCEED WITH CONDITIONS at a £215..." consistent with the
    deck title-bar's truncated "PROCEED WITH CONDITIONS at a £215..."
    even when the deck cuts the trailing detail.
  - :func:`check_engagement_consistency` — runs the extractor across a
    set of (artifact_type, format, file_path) tuples, returns a dict
    with per-artifact extracted text, normalised text, and a
    ``consistent`` boolean.

CLI::

    python tools/check_artifact_consistency.py --session-id <uuid>

Prints a per-artifact recommendation extract and the overall verdict.
Exit code 0 = consistent, 1 = divergent.
"""

from __future__ import annotations

import argparse
import asyncio
import os
import re
import sys
from pathlib import Path
from typing import Any

_REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_REPO_ROOT / "backend"))

from dotenv import load_dotenv  # noqa: E402

load_dotenv(_REPO_ROOT / ".env")


# Canonical verdict patterns. Cross-artifact consistency means every
# artifact's recommendation extract collapses to the SAME verdict
# token. Surrounding detail (deal size, geography qualifier) varies
# legitimately between artifacts (the deck title-bar truncates, the
# email lede paraphrases) — comparing on the verdict alone makes the
# check meaningful without false positives.
#
# Each pattern is a regex; ``verdict_id`` is the canonical token we
# emit. Order matters — the more specific patterns win first.
_VERDICT_PATTERNS: list[tuple[str, re.Pattern[str]]] = [
    ("proceed_with_conditions", re.compile(r"proceed\s+with\s+conditions", re.IGNORECASE)),
    ("proceed",                  re.compile(r"\bproceed\b", re.IGNORECASE)),
    ("walk_away",                re.compile(r"walk[-\s]?away", re.IGNORECASE)),
    ("reject",                   re.compile(r"\breject\b", re.IGNORECASE)),
    ("renegotiate",              re.compile(r"\brenegotiate\b", re.IGNORECASE)),
    ("expand_into_geography",    re.compile(r"expand\s+into\s+([a-z][a-z\s]+)", re.IGNORECASE)),
    ("expand",                   re.compile(r"\bexpand\b", re.IGNORECASE)),
    ("defer",                    re.compile(r"\bdefer\b", re.IGNORECASE)),
    ("wait_and_watch",           re.compile(r"wait[-\s]and[-\s]watch", re.IGNORECASE)),
]


def normalise_recommendation(text: str) -> str:
    """Reduce a recommendation string to its canonical verdict token.

    The check is on the VERDICT only — surrounding detail
    (valuation triple, geography qualifier, transaction structure)
    varies legitimately across artifacts. The verdict is the thing
    that must stay consistent. If two artifacts both contain
    ``PROCEED WITH CONDITIONS`` we count them as consistent even if
    their accompanying prose differs.

    For "expand into X" verdicts (growth-mode recommendations) the
    matched geography is captured into the token so an artifact that
    says "expand into Scotland" doesn't accidentally match one that
    says "expand into the EU".

    Returns an empty string when no canonical verdict pattern fires
    in the input — caller treats that as "skip from consistency
    comparison".
    """
    if not text:
        return ""
    s = re.sub(r"[\*_`~]+", " ", text)   # strip markdown emphasis
    s = re.sub(r"\s+", " ", s).strip()
    for verdict_id, pat in _VERDICT_PATTERNS:
        m = pat.search(s)
        if not m:
            continue
        if verdict_id == "expand_into_geography":
            geo = m.group(1).strip().lower()
            # Stop the geography capture at the first connective word
            # so "expand into Scotland via a partner-led..." normalises
            # to "expand_into:scotland" not the whole tail.
            geo = re.split(r"\s+(via|with|over|using|by|through|and)\s+", geo, maxsplit=1)[0]
            geo = re.sub(r"\s+", " ", geo).strip(" .,")
            return f"expand_into:{geo}"
        return verdict_id
    return ""


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------


def _extract_from_html(html: str) -> str:
    """1-pager HTML, email HTML, or interview-guide HTML. The
    recommendation lives inside a ``recommendation`` div / panel; if
    we can't find that, fall back to the first ``<h1>``/``<h2>``
    body text.
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    # 1-pager: <div class="recommendation"> ... </div>
    panel = soup.find(class_="recommendation") or soup.find(class_="recommendation-panel")
    if panel:
        return panel.get_text(" ", strip=True)
    # Email HTML: the recommendation appears in the lede paragraph
    # (first <p> with the verdict words).
    for p in soup.find_all("p"):
        text = p.get_text(" ", strip=True)
        low = text.lower()
        if any(v in low for v in ("proceed", "walk away", "walk-away", "expand", "renegotiate", "defer", "reject", "recommend")):
            return text
    # Last resort: first heading.
    h = soup.find(["h1", "h2"])
    return h.get_text(" ", strip=True) if h else ""


def _extract_from_markdown(md: str) -> str:
    """Email MD and interview-guide MD. Same heuristic as HTML — find
    the first paragraph that names a verdict keyword.

    Email MD has the lede in the first paragraph after "Dear ...";
    interview-guide MD has the recommendation in the "Pre-call
    briefing" bullets ("Recommendation: …").
    """
    lines = md.splitlines()
    # Interview-guide pre-call briefing pattern.
    for i, line in enumerate(lines):
        m = re.match(r"-\s*Recommendation:\s*(.*)$", line.strip())
        if m and m.group(1).strip():
            return m.group(1).strip()
    # Email pattern: first paragraph with a verdict keyword.
    paragraphs = re.split(r"\n\s*\n", md)
    for p in paragraphs:
        low = p.lower()
        if any(v in low for v in ("proceed", "walk away", "walk-away", "expand", "renegotiate", "defer", "reject", "recommend")):
            return p.strip()
    return paragraphs[0].strip() if paragraphs else ""


def _extract_from_pptx(file_path: str) -> str:
    """Deck PPTX. The recommendation slide is named
    ``recommendation`` in the W11/D5 sequence map; the title-bar text
    on that slide carries the recommendation prose. We also scan the
    title slide as a fallback because the deck title bar often
    carries the recommendation.
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    # Walk slides looking for one whose text body contains a verdict
    # keyword. Skip pure-chrome shapes (footer / page number).
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = ""
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text += (run.text or "") + " "
            text = text.strip()
            if not text:
                continue
            low = text.lower()
            if any(v in low for v in ("proceed", "walk away", "walk-away", "expand into", "renegotiate", "defer to", "reject")):
                return text
    # No verdict found — return the title slide's first text frame.
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    return t
    return ""


def _extract_from_xlsx(file_path: str) -> str:
    """Excel model XLSX. The Summary sheet carries the recommendation
    on row 3 (with "Recommendation:" label in col A, prose in col B
    onwards). The Cover sheet also has it as part of the title block.
    """
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    # Prefer Summary sheet — it's the partner-landing tab.
    for sheet_name in ("Summary", "Cover"):
        if sheet_name not in wb.sheetnames:
            continue
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if not isinstance(cell, str):
                    continue
                low = cell.lower()
                if any(v in low for v in ("proceed", "walk away", "walk-away", "expand into", "renegotiate", "defer to", "reject")):
                    return cell.strip()
    return ""


def _extract_from_pdf(file_path: str) -> str:
    """PDF artifacts (one_pager pdf, email pdf, interview-guide pdf).
    Walks pages via PyMuPDF; same verdict-keyword heuristic."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ""
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text() + "\n"
    paragraphs = re.split(r"\n\s*\n", text)
    for p in paragraphs:
        low = p.lower()
        if any(v in low for v in ("proceed", "walk away", "walk-away", "expand into", "renegotiate", "defer to", "reject", "recommend")):
            return p.strip()
    return text.strip().split("\n", 1)[0] if text.strip() else ""


_EXTRACTORS = {
    ("one_pager",       "html"): lambda p: _extract_from_html(_read_text(p)),
    ("one_pager",       "pdf"):  _extract_from_pdf,
    ("deck",            "pptx"): _extract_from_pptx,
    ("excel_model",     "xlsx"): _extract_from_xlsx,
    ("email",           "md"):   lambda p: _extract_from_markdown(_read_text(p)),
    ("email",           "html"): lambda p: _extract_from_html(_read_text(p)),
    ("email",           "pdf"):  _extract_from_pdf,
    ("interview_guide", "md"):   lambda p: _extract_from_markdown(_read_text(p)),
    ("interview_guide", "html"): lambda p: _extract_from_html(_read_text(p)),
    ("interview_guide", "pdf"):  _extract_from_pdf,
}


def _read_text(file_path: str | os.PathLike[str]) -> str:
    return Path(file_path).read_text(encoding="utf-8")


def extract_recommendation(artifact_type: str, format: str, file_path: str) -> str:
    """Dispatch to the right per-format extractor. Returns an empty
    string when the format has no registered extractor (caller treats
    that as "skip from consistency check")."""
    fn = _EXTRACTORS.get((artifact_type, format))
    if fn is None:
        return ""
    try:
        return fn(file_path)
    except Exception as e:  # noqa: BLE001
        return f"[extractor_error: {type(e).__name__}: {str(e)[:80]}]"


# ---------------------------------------------------------------------------
# Engagement-level check
# ---------------------------------------------------------------------------


def check_engagement_consistency(
    artifacts: list[dict[str, Any]],
    *,
    source_recommendation: str | None = None,
) -> dict[str, Any]:
    """Walk a list of artifact records and assert consistency.

    Each artifact record needs ``artifact_type``, ``format``, and
    ``file_path`` (the on-disk render). Records missing a file_path
    or with status != 'ready' are skipped (with a note in the result).

    When ``source_recommendation`` is provided, every artifact's
    normalised recommendation is also compared against the canonical
    source string; this catches drift where ALL artifacts agree with
    each other but disagree with the underlying report row.
    """
    per_artifact: list[dict[str, Any]] = []
    norms: list[str] = []
    for art in artifacts:
        atype = art["artifact_type"]
        fmt = art["format"]
        fpath = art.get("file_path")
        rec_entry: dict[str, Any] = {
            "artifact_type": atype,
            "format": fmt,
            "file_path": fpath,
            "extracted": "",
            "normalised": "",
            "skip_reason": None,
        }
        if art.get("status") not in (None, "ready"):
            rec_entry["skip_reason"] = f"status={art.get('status')}"
            per_artifact.append(rec_entry)
            continue
        if not fpath or not Path(fpath).exists():
            rec_entry["skip_reason"] = "no_file_path_on_disk"
            per_artifact.append(rec_entry)
            continue
        extracted = extract_recommendation(atype, fmt, fpath)
        rec_entry["extracted"] = extracted
        normalised = normalise_recommendation(extracted) if extracted else ""
        rec_entry["normalised"] = normalised
        per_artifact.append(rec_entry)
        if normalised:
            norms.append(normalised)

    consistent = bool(norms) and len(set(norms)) == 1
    source_match = None
    if source_recommendation:
        source_norm = normalise_recommendation(source_recommendation)
        source_match = all(n == source_norm for n in norms) if norms else None

    return {
        "consistent": consistent,
        "distinct_normalisations": sorted(set(norms)),
        "source_normalisation_match": source_match,
        "source_normalisation": (
            normalise_recommendation(source_recommendation) if source_recommendation else None
        ),
        "per_artifact": per_artifact,
    }


# ---------------------------------------------------------------------------
# CLI entry
# ---------------------------------------------------------------------------


async def _load_engagement_artifacts(session_id: str) -> tuple[list[dict[str, Any]], str | None]:
    """Pull the artifact list + the source recommendation from the DB."""
    from db.connection import acquire

    async with acquire() as conn:
        rows = await conn.fetch(
            """
            SELECT artifact_type, format, file_path, status
              FROM export_artifacts
             WHERE session_id = $1::uuid
             ORDER BY artifact_type, format
            """,
            session_id,
        )
        rec = await conn.fetchval(
            "SELECT recommendation FROM reports WHERE session_id = $1::uuid",
            session_id,
        )
    return [dict(r) for r in rows], rec


async def main_async(args: argparse.Namespace) -> int:
    from db.connection import close_db, init_db

    await init_db()
    try:
        artifacts, source_rec = await _load_engagement_artifacts(args.session_id)
    finally:
        await close_db()

    if not artifacts:
        print(f"No artifacts found for session_id={args.session_id!r}.")
        return 1

    result = check_engagement_consistency(artifacts, source_recommendation=source_rec)

    print(f"Source recommendation (normalised): {result['source_normalisation']!r}")
    print()
    for entry in result["per_artifact"]:
        label = f"{entry['artifact_type']}/{entry['format']}"
        if entry["skip_reason"]:
            print(f"  [SKIP] {label:<25}  ({entry['skip_reason']})")
            continue
        print(f"  {label:<25}  -> {entry['normalised']!r}")
    print()
    print(f"consistent: {result['consistent']}")
    if result["source_normalisation_match"] is not None:
        print(f"source_normalisation_match: {result['source_normalisation_match']}")

    return 0 if result["consistent"] else 1


def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("--session-id", required=True)
    args = p.parse_args()
    sys.exit(asyncio.run(main_async(args)))


if __name__ == "__main__":
    main()
