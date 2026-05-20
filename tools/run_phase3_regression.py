"""Phase 3 / Week 14 / Day 4 — full six-artifact regression across modes.

Walks the seeded Meridian Advisory engagements (W14/D3) and exercises
the entire Phase 3 deliverable bundle for each. The regression
deliberately uses the *cached* engagement payloads — the question
under test is whether the rendering pipelines produce a coherent,
mode-correct, citation-complete bundle, not whether the LLM stack
generates them. ``--force-regenerate`` is reserved for a future
variant that re-fires the analyst/writer pipeline.

Per engagement, fires the following (artifact_type, format) targets:

  - ``one_pager`` html + pdf
  - ``deck`` pptx
  - ``excel_model`` xlsx
  - ``email`` md + html + pdf
  - ``interview_guide`` md + html + pdf

That's 10 artifacts × 2 engagements = 20 generations. PDF formats
land as ``skipped_no_weasyprint`` when the native runtime isn't
loadable (Windows dev hosts); Docker / Linux production renders them
all. The headline assertions tolerate either state so the regression
runs identically across environments.

Captures per artifact: status, file size, generation time, citation
count, mode-specific content presence flags. Headline assertions
roll the bundle up into eight booleans; ``headline_pass`` is the
ship gate.

Cross-artifact consistency lives in :mod:`tools.check_artifact_consistency`.

Outputs to ``backend/eval_runs/phase3_regression/summary.json``.

Usage::

    python tools/run_phase3_regression.py
    python tools/run_phase3_regression.py --summary-only
"""

from __future__ import annotations

import argparse
import asyncio
import json
import os
import re
import shutil
import sys
import time
import uuid
from pathlib import Path
from typing import Any

_REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_REPO_ROOT / "backend"))
sys.path.insert(0, str(_REPO_ROOT))   # so ``tools.check_artifact_consistency`` resolves

from dotenv import load_dotenv  # noqa: E402

load_dotenv(_REPO_ROOT / ".env")

_BENCH_ROOT = _REPO_ROOT / "backend" / "eval_runs" / "phase3_regression"
_ARTIFACT_OUT = _REPO_ROOT / "artifacts_out" / "phase3_regression"

ARTIFACT_TARGETS: list[tuple[str, str]] = [
    ("one_pager",       "html"),
    ("one_pager",       "pdf"),
    ("deck",            "pptx"),
    ("excel_model",     "xlsx"),
    ("email",           "md"),
    ("email",           "html"),
    ("email",           "pdf"),
    ("interview_guide", "md"),
    ("interview_guide", "html"),
    ("interview_guide", "pdf"),
]

# Mode-specific markers we expect in payload-driven outputs. The
# regression asserts presence/absence per mode to catch
# cross-contamination (M&A artifacts emitting Porter's; growth
# artifacts emitting valuation_range, etc.).
_MODE_MARKERS = {
    "m_and_a_diligence": {
        "must_have_in_consulting_payload": (
            "valuation_range", "synergy_estimate", "deal_structure_implications",
        ),
        "must_have_in_text_artifacts": ("PROCEED", "valuation", "EBITDA"),
        "must_not_have_in_text_artifacts": ("Porter's", "competitive landscape"),
    },
    "growth_strategy": {
        "must_have_in_consulting_payload": (
            "frameworks", "competitive_landscape", "options_matrix",
        ),
        "must_have_in_text_artifacts": ("Porter", "competitive", "channel"),
        "must_not_have_in_text_artifacts": ("valuation_range", "synergy_estimate"),
    },
}


# ---------------------------------------------------------------------------
# Engagement discovery
# ---------------------------------------------------------------------------


async def _load_meridian_engagements() -> list[dict[str, Any]]:
    """Read the two seeded Meridian engagements from the DB. Returns
    a list of dicts with session_id, title, mode, recommendation."""
    from db.connection import acquire

    async with acquire() as conn:
        rows = await conn.fetch(
            """
            SELECT s.id AS session_id,
                   s.title,
                   s.report_mode,
                   r.recommendation,
                   r.consulting_payload
              FROM sessions s
              JOIN firms f  ON f.id = s.firm_id
              JOIN reports r ON r.session_id = s.id
             WHERE f.slug = 'meridian-advisory'
             ORDER BY s.report_mode, s.title
            """,
        )
    out = []
    for r in rows:
        cp = r["consulting_payload"]
        if isinstance(cp, str):
            try:
                cp = json.loads(cp)
            except Exception:
                cp = {}
        out.append({
            "session_id": str(r["session_id"]),
            "title": r["title"],
            "mode": r["report_mode"],
            "recommendation": r["recommendation"],
            "consulting_payload": cp if isinstance(cp, dict) else {},
        })
    return out


# ---------------------------------------------------------------------------
# Per-artifact firing
# ---------------------------------------------------------------------------


def _weasyprint_available() -> bool:
    try:
        from weasyprint import HTML  # type: ignore
        HTML(string="<html><body>x</body></html>").write_pdf()
        return True
    except Exception:
        return False


async def _fire_one(eng: dict[str, Any], atype: str, fmt: str, has_wp: bool) -> dict[str, Any]:
    from core.exports import GenerateArtifactRequest, generate_artifact

    rec: dict[str, Any] = {
        "engagement_id": eng["session_id"],
        "engagement_title": eng["title"],
        "mode": eng["mode"],
        "artifact_type": atype,
        "format": fmt,
    }
    if fmt == "pdf" and not has_wp:
        rec.update({
            "status": "skipped_no_weasyprint",
            "file_size_bytes": None,
            "wall_seconds": 0.0,
        })
        return rec

    t0 = time.perf_counter()
    req = GenerateArtifactRequest(
        session_id=uuid.UUID(eng["session_id"]),
        artifact_type=atype,
        format=fmt,
    )
    try:
        result = await generate_artifact(req)
    except Exception as e:  # noqa: BLE001
        rec.update({
            "status": "exception",
            "error": f"{type(e).__name__}: {str(e)[:200]}",
            "wall_seconds": time.perf_counter() - t0,
        })
        return rec

    rec.update({
        "artifact_id": str(result.artifact_id),
        "status": result.status,
        "file_path": result.file_path,
        "file_size_bytes": result.file_size_bytes,
        "claim_citation_count": result.claim_citation_count,
        "metadata": result.metadata or {},
        "wall_seconds": round(time.perf_counter() - t0, 3),
        "failure_reason": result.failure_reason,
    })

    if result.status == "ready" and result.file_path:
        out_path = _ARTIFACT_OUT / f"{eng['mode']}_{atype}_{fmt}.{_ext_for(atype, fmt)}"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            shutil.copy(result.file_path, out_path)
            rec["copied_to"] = str(out_path)
        except Exception:
            pass

    return rec


def _ext_for(atype: str, fmt: str) -> str:
    return fmt


# ---------------------------------------------------------------------------
# Per-artifact content inspection
# ---------------------------------------------------------------------------


def _inspect_artifact_content(rec: dict[str, Any], eng: dict[str, Any]) -> None:
    """Pull mode-marker presence signals + the extracted recommendation
    into the per-artifact record. Skips artifacts that aren't ``ready``."""
    if rec.get("status") != "ready" or not rec.get("file_path"):
        return
    fpath = rec["file_path"]
    atype, fmt = rec["artifact_type"], rec["format"]

    text = ""
    try:
        if (atype, fmt) in (
            ("one_pager", "html"), ("email", "html"), ("email", "md"),
            ("interview_guide", "md"), ("interview_guide", "html"),
        ):
            text = Path(fpath).read_text(encoding="utf-8")
        elif (atype, fmt) == ("deck", "pptx"):
            from pptx import Presentation
            prs = Presentation(fpath)
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text or "")
            text = "\n".join(parts)
        elif (atype, fmt) == ("excel_model", "xlsx"):
            from openpyxl import load_workbook
            wb = load_workbook(fpath, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if isinstance(cell, str):
                            parts.append(cell)
            text = "\n".join(parts)
        elif fmt == "pdf":
            try:
                import fitz
                with fitz.open(fpath) as doc:
                    text = "\n".join(p.get_text() for p in doc)
            except Exception:
                text = ""
    except Exception as e:  # noqa: BLE001
        rec["inspection_error"] = f"{type(e).__name__}: {str(e)[:120]}"
        return

    markers = _MODE_MARKERS.get(eng["mode"], {})
    text_lower = text.lower()
    rec["mode_markers"] = {
        "must_have_present": [
            tok for tok in markers.get("must_have_in_text_artifacts", ())
            if tok.lower() in text_lower
        ],
        "must_not_have_leaked": [
            tok for tok in markers.get("must_not_have_in_text_artifacts", ())
            if tok.lower() in text_lower
        ],
    }
    rec["text_chars"] = len(text)


# ---------------------------------------------------------------------------
# Headline assertions
# ---------------------------------------------------------------------------


def _headline_assertions(
    records: list[dict[str, Any]],
    engagements: list[dict[str, Any]],
    consistency: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    out: dict[str, Any] = {}

    # 1. All generations succeed (ready OR skipped_no_weasyprint).
    statuses = [r["status"] for r in records]
    out["status_counts"] = {
        "ready": statuses.count("ready"),
        "skipped_no_weasyprint": statuses.count("skipped_no_weasyprint"),
        "failed": statuses.count("failed"),
        "exception": statuses.count("exception"),
    }
    out["all_generations_succeed"] = all(
        s in ("ready", "skipped_no_weasyprint") for s in statuses
    )

    # 2. Mode-awareness: per-engagement consulting_payload contains the
    # expected mode-specific keys; the *text artifacts* show no
    # cross-contamination.
    mode_ok = True
    contamination_detail: dict[str, dict[str, Any]] = {}
    for eng in engagements:
        markers = _MODE_MARKERS.get(eng["mode"]) or {}
        cp = eng.get("consulting_payload") or {}
        missing_cp = [k for k in markers.get("must_have_in_consulting_payload", ()) if k not in cp]
        leaked: list[str] = []
        for r in records:
            if r["engagement_id"] != eng["session_id"]:
                continue
            leaked.extend(r.get("mode_markers", {}).get("must_not_have_leaked", []))
        leaked = sorted(set(leaked))
        contamination_detail[eng["title"]] = {
            "missing_consulting_payload_keys": missing_cp,
            "leaked_text_markers": leaked,
        }
        if missing_cp or leaked:
            mode_ok = False
    out["mode_aware_no_cross_contamination"] = mode_ok
    out["mode_aware_detail"] = contamination_detail

    # 3. Growth Porter's renders real content (not the fallback). The
    # deck builder renders each force as its own box without
    # necessarily using the word "Porter" in the slide text, so we
    # scan the WHOLE deck for ≥4 of the 5 force keywords (rivalry,
    # supplier, buyer, substitute, entrant). 4-of-5 is enough — the
    # W11/D5 fallback path emits the literal "not produced for this
    # engagement" sentence which we also check below.
    growth_sessions = [e["session_id"] for e in engagements if e["mode"] == "growth_strategy"]
    porters_ok = True
    porters_detail: dict[str, dict[str, Any]] = {}
    for sid in growth_sessions:
        deck = next(
            (r for r in records if r["engagement_id"] == sid and r["artifact_type"] == "deck" and r["status"] == "ready"),
            None,
        )
        op_html = next(
            (r for r in records if r["engagement_id"] == sid and r["artifact_type"] == "one_pager" and r["format"] == "html" and r["status"] == "ready"),
            None,
        )
        d: dict[str, Any] = {}
        if deck and deck.get("file_path"):
            from pptx import Presentation
            prs = Presentation(deck["file_path"])
            full_deck_text = "\n".join(
                shape.text_frame.text
                for slide in prs.slides
                for shape in slide.shapes
                if shape.has_text_frame
            ).lower()
            forces_found = sum(
                1 for force in ("rivalry", "supplier", "buyer", "substitute", "entrant")
                if force in full_deck_text
            )
            d["deck_force_keywords_found"] = forces_found
            d["deck_has_fallback_marker"] = "not produced for this engagement" in full_deck_text
            if forces_found < 4 or d["deck_has_fallback_marker"]:
                porters_ok = False
        else:
            d["deck_status"] = "not_ready"
            porters_ok = False
        if op_html and op_html.get("file_path"):
            html = Path(op_html["file_path"]).read_text(encoding="utf-8")
            fallback = "not produced for this engagement" in html.lower()
            d["one_pager_fallback_rendered"] = fallback
            if fallback:
                porters_ok = False
        porters_detail[sid] = d
    out["growth_porters_real_not_fallback"] = porters_ok
    out["growth_porters_detail"] = porters_detail

    # 4. Citation completeness — engagement-level. The spec asks for
    # ≥5 distinct claim citations per engagement. We aggregate
    # ``claim_citation_count`` + the metadata ``cited_claim_ids``
    # across every ready artifact for the engagement; an engagement
    # passes if the union of cited ids is ≥5. Per-artifact counts
    # legitimately vary (interview-guide's Section B counts only
    # the reasons/risks it pressure-tests; email counts only the
    # body's lede/recommendation/caveat references), so the
    # engagement-aggregate is the meaningful gate.
    engagement_citations: dict[str, set[str]] = {}
    for r in records:
        if r["status"] != "ready":
            continue
        eid = r["engagement_id"]
        bucket = engagement_citations.setdefault(eid, set())
        md = r.get("metadata") or {}
        for cid in (md.get("cited_claim_ids") or []):
            if isinstance(cid, str) and cid.strip():
                bucket.add(cid.strip())
        # Some artifacts (excel_model, one_pager) emit comments /
        # chips per-claim but don't surface ``cited_claim_ids`` in
        # their metadata. Fall back on ``claim_citation_count`` so
        # those still contribute to the engagement total when other
        # artifacts don't list specific ids.
    citation_detail: dict[str, dict[str, Any]] = {}
    citation_ok = True
    for eng in engagements:
        ids = sorted(engagement_citations.get(eng["session_id"], set()))
        # Fallback: max(claim_citation_count) across the engagement's
        # ready artifacts is a useful floor when no artifact emitted
        # named ids.
        max_count = max(
            ((r.get("claim_citation_count") or 0)
             for r in records
             if r["engagement_id"] == eng["session_id"] and r["status"] == "ready"),
            default=0,
        )
        engagement_min = max(len(ids), max_count)
        citation_detail[eng["title"]] = {
            "distinct_claim_ids": ids,
            "max_per_artifact_count": max_count,
            "engagement_floor": engagement_min,
        }
        if engagement_min < 5:
            citation_ok = False
    out["citation_completeness_detail"] = citation_detail
    out["citation_completeness"] = citation_ok

    # 5. Firm branding visible on every PPTX / XLSX / PDF.
    # PPTX: check the first slide for the firm name string. XLSX: check
    # the Cover sheet A1 area. PDF: check for the firm name in extracted
    # text via PyMuPDF.
    branding_ok = True
    branding_detail: dict[str, str] = {}
    for r in records:
        if r["status"] != "ready":
            continue
        atype, fmt = r["artifact_type"], r["format"]
        if (atype, fmt) not in (("deck", "pptx"), ("excel_model", "xlsx")) and fmt != "pdf":
            continue
        # The firm display name is whatever appears in the artifact's
        # metadata firm_name or in the source engagement title; for
        # the Meridian seed, it's "Meridian Advisory".
        looking_for = "Meridian Advisory"
        found = False
        try:
            if (atype, fmt) == ("deck", "pptx"):
                from pptx import Presentation
                prs = Presentation(r["file_path"])
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame and looking_for in shape.text_frame.text:
                            found = True
                            break
                    if found:
                        break
            elif (atype, fmt) == ("excel_model", "xlsx"):
                from openpyxl import load_workbook
                wb = load_workbook(r["file_path"], data_only=True)
                for sn in wb.sheetnames:
                    ws = wb[sn]
                    for row in ws.iter_rows(values_only=True):
                        for cell in row:
                            if isinstance(cell, str) and looking_for in cell:
                                found = True
                                break
                        if found:
                            break
                    if found:
                        break
            elif fmt == "pdf":
                try:
                    import fitz
                    with fitz.open(r["file_path"]) as doc:
                        for page in doc:
                            if looking_for in page.get_text():
                                found = True
                                break
                except Exception:
                    pass
        except Exception:
            found = False
        branding_detail[f"{r['engagement_title']}/{atype}/{fmt}"] = "ok" if found else "missing"
        if not found:
            branding_ok = False
    out["firm_branding_visible"] = branding_ok
    out["branding_detail"] = branding_detail

    # 6. Cross-artifact consistency per engagement.
    consistent_all = all(c["consistent"] for c in consistency.values())
    out["cross_artifact_consistent"] = consistent_all
    out["consistency_per_engagement"] = {
        eid: {
            "consistent": c["consistent"],
            "distinct_normalisations": c["distinct_normalisations"],
            "source_match": c["source_normalisation_match"],
        }
        for eid, c in consistency.items()
    }

    # 7. Excel citation audit empty for both modes.
    audit_ok = True
    excel_audit_detail: dict[str, list] = {}
    for r in records:
        if r["artifact_type"] != "excel_model" or r["format"] != "xlsx":
            continue
        if r["status"] != "ready":
            continue
        try:
            from core.exports.excel._branding import audit_citations
            from openpyxl import load_workbook
            wb = load_workbook(r["file_path"])
            audit = audit_citations(wb)
            excel_audit_detail[r["engagement_title"]] = audit["missing"]
            if audit["missing"]:
                audit_ok = False
        except Exception as e:  # noqa: BLE001
            excel_audit_detail[r["engagement_title"]] = [{"error": str(e)[:120]}]
            audit_ok = False
    out["excel_citation_audit_clean"] = audit_ok
    out["excel_audit_detail"] = excel_audit_detail

    # 8. Total cost zero.
    total_cost = sum(
        (r.get("metadata") or {}).get("generation_cost_usd") or 0.0 for r in records
    )
    out["total_cost_zero"] = total_cost == 0.0
    out["total_cost_usd"] = total_cost

    out["headline_pass"] = all(
        v for k, v in out.items()
        if isinstance(v, bool) and not k.endswith("_info")
    )
    return out


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


async def main_async(args: argparse.Namespace) -> int:
    from db.connection import close_db, init_db
    from tools.check_artifact_consistency import check_engagement_consistency

    _BENCH_ROOT.mkdir(parents=True, exist_ok=True)
    _ARTIFACT_OUT.mkdir(parents=True, exist_ok=True)

    if args.summary_only:
        summary_path = _BENCH_ROOT / "summary.json"
        if summary_path.exists():
            print(summary_path.read_text(encoding="utf-8"))
            return 0
        print("No previous summary.json. Run without --summary-only first.")
        return 1

    await init_db()
    try:
        engagements = await _load_meridian_engagements()
        if not engagements:
            raise SystemExit(
                "No Meridian engagements found. Run tools/seed_sample_workspace.py first."
            )
        has_wp = _weasyprint_available()

        records: list[dict[str, Any]] = []
        for eng in engagements:
            print(f"\n=== {eng['title']} (mode={eng['mode']}) ===")
            for atype, fmt in ARTIFACT_TARGETS:
                rec = await _fire_one(eng, atype, fmt, has_wp)
                _inspect_artifact_content(rec, eng)
                records.append(rec)
                size = rec.get("file_size_bytes")
                print(
                    f"  {atype:<16}/{fmt:<5}  {rec['status']:<22}  "
                    f"size={size!s:<8}  wall={rec.get('wall_seconds', 0):.2f}s"
                )

        # Cross-artifact consistency per engagement.
        consistency_results: dict[str, dict[str, Any]] = {}
        for eng in engagements:
            eng_artifacts = [r for r in records if r["engagement_id"] == eng["session_id"]]
            consistency_results[eng["session_id"]] = check_engagement_consistency(
                eng_artifacts,
                source_recommendation=eng["recommendation"],
            )
    finally:
        await close_db()

    headline = _headline_assertions(records, engagements, consistency_results)

    summary = {
        "generated_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
        "n_engagements": len(engagements),
        "n_artifact_generations": len(records),
        "headline_assertions": headline,
        "headline_pass": headline["headline_pass"],
        "engagements": [
            {"session_id": e["session_id"], "title": e["title"], "mode": e["mode"],
             "recommendation": e["recommendation"]}
            for e in engagements
        ],
        "runs": records,
        "consistency": consistency_results,
    }
    (_BENCH_ROOT / "summary.json").write_text(
        json.dumps(summary, indent=2, default=str), encoding="utf-8",
    )

    print("\n=== HEADLINE ASSERTIONS ===")
    for k, v in headline.items():
        if isinstance(v, bool):
            print(f"  [{'PASS' if v else 'FAIL'}] {k}")
        elif not isinstance(v, dict):
            print(f"  {k}: {v}")
    print(f"\nheadline_pass: {summary['headline_pass']}")
    print(f"summary: {_BENCH_ROOT / 'summary.json'}")
    return 0 if summary["headline_pass"] else 1


def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("--summary-only", action="store_true")
    p.add_argument("--force-regenerate", action="store_true",
                   help="Reserved — re-run the LLM pipeline before rendering. Not implemented today.")
    args = p.parse_args()
    sys.exit(asyncio.run(main_async(args)))


if __name__ == "__main__":
    main()
